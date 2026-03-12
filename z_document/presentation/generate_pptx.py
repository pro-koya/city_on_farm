#!/usr/bin/env python3
"""
セッツマルシェ プレゼンテーション PPTX 生成スクリプト
Google Slides にインポート可能な .pptx を生成します。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DIR = os.path.dirname(os.path.abspath(__file__))
IMG = lambda name: os.path.join(DIR, name)

MAIN      = RGBColor(0x4C, 0x6B, 0x5C)
MAIN_DARK = RGBColor(0x35, 0x5D, 0x4A)
ACCENT    = RGBColor(0xA3, 0xC9, 0xA8)
HIGHLIGHT = RGBColor(0xD6, 0xEA, 0xDF)
SOFT      = RGBColor(0xF2, 0xF7, 0xF4)
BG        = RGBColor(0xF9, 0xF9, 0xF6)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_CLR  = RGBColor(0x2B, 0x2B, 0x2B)
MUTED     = RGBColor(0x6B, 0x72, 0x80)
WARN_BG   = RGBColor(0xE8, 0x92, 0x7C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_JP = 'Noto Sans JP'
FONT_SERIF = 'Noto Serif JP'

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank

# ─── helpers ───────────────────────────────────────────────────
def add_bg_color(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_bg_gradient(slide, c1, c2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0

def add_image_bg(slide, img_path, alpha=0):
    slide.shapes.add_picture(img_path, Emu(0), Emu(0), SLIDE_W, SLIDE_H)

def add_overlay(slide, color, alpha_pct=40):
    from lxml import etree
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    sp_elem = shp._element
    solid_fill = sp_elem.find('.//' + qn('a:solidFill'))
    if solid_fill is not None:
        srgb = solid_fill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_el = etree.SubElement(srgb, qn('a:alpha'))
            alpha_el.set('val', str(alpha_pct * 1000))
    shp.line.fill.background()

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_CLR, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_JP, anchor=MSO_ANCHOR.TOP, line_spacing=1.4):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    try:
        p.line_spacing = Pt(int(font_size * line_spacing))
    except:
        pass
    return txBox

def add_multi_text(slide, left, top, width, height, runs_list,
                   alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs_list = [(text, size, color, bold, font_name), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, size, color, bold, font_name) in enumerate(runs_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name or FONT_JP
        try:
            p.line_spacing = Pt(int(size * 1.5))
        except:
            pass
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color,
                     border_color=None, border_width=Pt(0)):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = border_width
    else:
        shp.line.fill.background()
    adj = shp.adjustments
    if len(adj) > 0:
        adj[0] = 0.05
    return shp

def add_stripe_bottom(slide):
    h = Inches(0.06)
    w_third = Emu(SLIDE_W.emu // 3)
    for i, c in enumerate([MAIN, ACCENT, HIGHLIGHT]):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(w_third * i),
            Emu(SLIDE_H.emu - h.emu),
            w_third + Emu(10),
            h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = c
        shp.line.fill.background()

def add_slide_num(slide, num, total):
    add_textbox(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.35),
                f'{num} / {total}', font_size=9, color=MUTED,
                alignment=PP_ALIGN.RIGHT)

def add_screenshot_placeholder(slide, left, top, width, height, label):
    shp = add_rounded_rect(slide, left, top, width, height, SOFT, ACCENT, Pt(1.5))
    shp.line.dash_style = 4  # dash
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0)
    r = p.add_run()
    r.text = '\n\n'
    r.font.size = Pt(8)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(11)
    r2.font.color.rgb = MUTED
    r2.font.name = FONT_JP
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = '[ スクリーンショットを挿入 ]'
    r3.font.size = Pt(10)
    r3.font.color.rgb = ACCENT
    r3.font.name = FONT_JP
    r3.font.bold = True

def add_card(slide, left, top, width, height, title, body,
             accent_color=MAIN, bg_color=WHITE):
    card = add_rounded_rect(slide, left, top, width, height, bg_color, RGBColor(0xE5,0xE7,0xEB), Pt(0.75))
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    add_textbox(slide, Emu(left.emu + Inches(0.25).emu), Emu(top.emu + Inches(0.2).emu),
                Emu(width.emu - Inches(0.4).emu), Inches(0.4),
                title, font_size=16, bold=True, color=MAIN_DARK)
    add_textbox(slide, Emu(left.emu + Inches(0.25).emu), Emu(top.emu + Inches(0.65).emu),
                Emu(width.emu - Inches(0.4).emu), Emu(height.emu - Inches(0.85).emu),
                body, font_size=12, color=MUTED, line_spacing=1.6)

TOTAL = 20

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — タイトル
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_image_bg(s, IMG('slide_hero.png'))
add_overlay(s, MAIN_DARK, alpha_pct=55)
add_textbox(s, Inches(0.8), Inches(1.2), Inches(5), Inches(0.5),
            '2026.03.24  ご紹介資料', font_size=13, color=WHITE,
            bold=False, alignment=PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(2.0), Inches(8), Inches(1.2),
            'セッツマルシェ', font_size=54, color=WHITE,
            bold=True, font_name=FONT_SERIF, alignment=PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(3.5), Inches(8), Inches(0.8),
            'まちと畑を、ひとつの食卓につなぐ。', font_size=24, color=WHITE,
            bold=False, font_name=FONT_SERIF, alignment=PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(5.8), Inches(6), Inches(0.5),
            '渡辺ファーム  ／  摂津市商工会', font_size=12, color=RGBColor(0xCC,0xCC,0xCC),
            alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — アジェンダ
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
add_textbox(s, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            '本日の内容', font_size=32, bold=True, color=MAIN_DARK,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER)

agenda_items = [
    ('1', 'セッツマルシェとは', '背景・コンセプト・目指す姿'),
    ('2', 'わたしたちの想い', '小さな声を拾う ／ まちの食を支える ／ 本当の美味しさを届ける'),
    ('3', '取り扱う野菜の出どころ', '渡辺ファーム ／ 小規模農家 ／ 市民農園'),
    ('4', '購入フローのご紹介', '購入手順・決済方法・配送・畑受け取り'),
    ('5', '便利な機能と今後', 'リピート注文・テンプレート・掛売・法人向け機能'),
]
for i, (num, title, sub) in enumerate(agenda_items):
    y = Inches(1.8 + i * 1.0)
    circ = slide_shapes = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(3.5), y, Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = MAIN
    circ.line.fill.background()
    tf = circ.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    r.font.size = Pt(16)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = FONT_JP
    add_textbox(s, Inches(4.2), y, Inches(5.5), Inches(0.3),
                title, font_size=17, bold=True, color=MAIN_DARK)
    add_textbox(s, Inches(4.2), Emu(y.emu + Inches(0.35).emu), Inches(5.5), Inches(0.3),
                sub, font_size=11, color=MUTED)

add_stripe_bottom(s)
add_slide_num(s, 2, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — 背景・課題
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
add_textbox(s, Inches(1), Inches(0.5), Inches(11.3), Inches(0.9),
            '地域の「もったいない」を、\n「おいしい」に変えたい',
            font_size=30, bold=True, color=MAIN_DARK,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

s.shapes.add_picture(IMG('slide_challenges.png'),
                     Inches(0.8), Inches(1.9), Inches(11.7), Inches(4.2))

challenges = [
    ('小規模農家の課題', '採れすぎた野菜が販路に\n乗らず廃棄されてしまう'),
    ('飲食店の課題', '新鮮な地場野菜を安定的に\n仕入れるルートが限られている'),
    ('地域の課題', 'つくる人・届ける人・味わう人の\nつながりが見えにくくなっている'),
]
for i, (t, d) in enumerate(challenges):
    x = Inches(1.0 + i * 4.0)
    add_textbox(s, x, Inches(6.15), Inches(3.5), Inches(0.3),
                t, font_size=13, bold=True, color=MAIN_DARK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, Inches(6.5), Inches(3.5), Inches(0.6),
                d, font_size=10, color=MUTED,
                alignment=PP_ALIGN.CENTER, line_spacing=1.5)

add_stripe_bottom(s)
add_slide_num(s, 3, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — セッツマルシェとは
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, SOFT)
add_textbox(s, Inches(1), Inches(0.5), Inches(11.3), Inches(0.6),
            'セッツマルシェとは', font_size=32, bold=True, color=MAIN_DARK,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(2), Inches(1.2), Inches(9.3), Inches(0.5),
            '市民農園・小規模農家の野菜を、地域の飲食店へ届けるオンラインマルシェ',
            font_size=15, color=MUTED, alignment=PP_ALIGN.CENTER)

box = add_rounded_rect(s, Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.8), WHITE, ACCENT, Pt(1))
add_textbox(s, Inches(1.8), Inches(2.3), Inches(9.7), Inches(4.0),
            '食材を育てる人と、それを届ける人、味わう人。\n'
            'そのすべてがつながることで、本当のおいしさは生まれます。\n\n'
            '私たちは、市民農園や小規模農家さんが心を込めて育てた野菜を、\n'
            '地域の飲食店へと届ける仕組みを作りました。\n\n'
            '「少し作り過ぎてしまった野菜」も、誰かの料理で光を放つ。\n'
            'そんな食の循環が、地域をもっと豊かにしていくと信じています。',
            font_size=15, color=TEXT_CLR, alignment=PP_ALIGN.CENTER,
            line_spacing=1.8, font_name=FONT_SERIF)

add_stripe_bottom(s)
add_slide_num(s, 4, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — 3つの柱
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
add_textbox(s, Inches(1), Inches(0.4), Inches(11.3), Inches(0.6),
            'わたしたちの想い', font_size=32, bold=True, color=MAIN_DARK,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER)

s.shapes.add_picture(IMG('slide_pillars.png'),
                     Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.3))

pillars = [
    ('小さな声を拾う', '家庭菜園や小規模農家の収穫を\n活かし、販売機会に変える', MAIN),
    ('まちの食を支える', '地域の飲食店に、新鮮で確かな\n食材を安定して届ける', ACCENT),
    ('本当の美味しさを届ける', '採れたての鮮度と生産者の想いを\nそのまま、食卓へつなげる', HIGHLIGHT),
]
for i, (t, d, c) in enumerate(pillars):
    x = Inches(0.8 + i * 4.1)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(5.75), Inches(3.5), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = c
    bar.line.fill.background()
    add_textbox(s, x, Inches(5.95), Inches(3.5), Inches(0.4),
                t, font_size=16, bold=True, color=MAIN_DARK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, Inches(6.4), Inches(3.5), Inches(0.7),
                d, font_size=11, color=MUTED,
                alignment=PP_ALIGN.CENTER, line_spacing=1.6)

add_stripe_bottom(s)
add_slide_num(s, 5, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — 野菜の出どころ
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
add_textbox(s, Inches(1), Inches(0.4), Inches(11.3), Inches(0.6),
            '取り扱う野菜の出どころ', font_size=32, bold=True, color=MAIN_DARK,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER)

s.shapes.add_picture(IMG('slide_sources.png'),
                     Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.3))

sources = [
    ('市民農園の方々', '家庭菜園の余剰も、飲食店の\n一皿に。小さな生産にも価値を。', HIGHLIGHT),
    ('小規模農家さん', '地域で丁寧に育てた野菜。\n販路を共有し、収穫を無駄にしない。', ACCENT),
    ('渡辺ファーム', '運営元の農場。安心・安全な栽培で、\n主力となる野菜を安定供給。', MAIN),
]
for i, (t, d, c) in enumerate(sources):
    x = Inches(0.8 + i * 4.1)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(5.75), Inches(3.5), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = c
    bar.line.fill.background()
    add_textbox(s, x, Inches(5.95), Inches(3.5), Inches(0.4),
                t, font_size=16, bold=True, color=MAIN_DARK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, Inches(6.4), Inches(3.5), Inches(0.7),
                d, font_size=11, color=MUTED,
                alignment=PP_ALIGN.CENTER, line_spacing=1.6)

add_stripe_bottom(s)
add_slide_num(s, 6, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — 目指す姿
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_image_bg(s, IMG('slide_vision.png'))
add_overlay(s, MAIN_DARK, alpha_pct=50)

add_textbox(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(2.5),
            '地域でつくり、\n地域で味わい、\n地域でつながる。',
            font_size=40, bold=True, color=WHITE,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER, line_spacing=1.5)
add_textbox(s, Inches(2.5), Inches(4.5), Inches(8.3), Inches(1.5),
            'この循環の輪が、まちの豊かさを育んでいく。\nそれが、私たち「セッツマルシェ」の願いです。',
            font_size=17, color=RGBColor(0xDD,0xDD,0xDD),
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER, line_spacing=1.7)
add_slide_num(s, 7, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Part 2 セクション表紙
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, SOFT)

label = add_rounded_rect(s, Inches(5.4), Inches(2.0), Inches(2.5), Inches(0.5), MAIN)
tf = label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = 'PART 2'
r.font.size = Pt(13)
r.font.color.rgb = WHITE
r.font.bold = True
r.font.name = FONT_JP

add_textbox(s, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.0),
            'ご利用方法のご紹介', font_size=36, bold=True, color=MAIN_DARK,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(2.5), Inches(4.2), Inches(8.3), Inches(0.5),
            '飲食店さま向け — 購入フローと便利な機能', font_size=15, color=MUTED,
            alignment=PP_ALIGN.CENTER)
add_stripe_bottom(s)
add_slide_num(s, 8, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — 購入の流れ 全体像
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
add_textbox(s, Inches(1), Inches(0.4), Inches(11.3), Inches(0.6),
            '購入の流れ', font_size=32, bold=True, color=MAIN_DARK,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER)

s.shapes.add_picture(IMG('slide_flow_overview.png'),
                     Inches(1.0), Inches(1.4), Inches(11.3), Inches(4.2))

flow_labels = ['商品を探す', 'カートに入れる', '購入手続き', '決済', '注文確定']
for i, label in enumerate(flow_labels):
    x = Inches(1.3 + i * 2.25)
    add_textbox(s, x, Inches(5.7), Inches(1.8), Inches(0.4),
                label, font_size=13, bold=True, color=MAIN_DARK,
                alignment=PP_ALIGN.CENTER)

tags = ['クレジットカード', '代金引換', '掛売（請求書払い）', '配送', '畑で受け取り']
for i, t in enumerate(tags):
    x = Inches(1.8 + i * 2.0)
    tag_shp = add_rounded_rect(s, x, Inches(6.4), Inches(1.7), Inches(0.35), HIGHLIGHT, ACCENT, Pt(0.5))
    tf = tag_shp.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = t
    r.font.size = Pt(10)
    r.font.color.rgb = MAIN_DARK
    r.font.name = FONT_JP

add_stripe_bottom(s)
add_slide_num(s, 9, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — STEP 1: 商品を探す
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
step_label = add_rounded_rect(s, Inches(0.8), Inches(0.5), Inches(1.2), Inches(0.4), MAIN)
tf = step_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = 'STEP 1'
r.font.size = Pt(11)
r.font.color.rgb = WHITE
r.font.bold = True
r.font.name = FONT_JP

add_textbox(s, Inches(2.2), Inches(0.45), Inches(6), Inches(0.5),
            '商品を探す', font_size=26, bold=True, color=MAIN_DARK)
add_textbox(s, Inches(0.8), Inches(1.1), Inches(6), Inches(0.3),
            'トップページや商品一覧から、旬の野菜を探せます',
            font_size=13, color=MUTED)

add_screenshot_placeholder(s, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.2),
                           '商品一覧ページ')
add_screenshot_placeholder(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.2),
                           '商品詳細ページ')

add_stripe_bottom(s)
add_slide_num(s, 10, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — STEP 2: カートに入れる
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
step_label = add_rounded_rect(s, Inches(0.8), Inches(0.5), Inches(1.2), Inches(0.4), MAIN)
tf = step_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = 'STEP 2'
r.font.size = Pt(11)
r.font.color.rgb = WHITE
r.font.bold = True
r.font.name = FONT_JP

add_textbox(s, Inches(2.2), Inches(0.45), Inches(6), Inches(0.5),
            'カートに入れる', font_size=26, bold=True, color=MAIN_DARK)
add_textbox(s, Inches(0.8), Inches(1.1), Inches(6), Inches(0.3),
            '数量を選んで「カートに追加」。複数の商品をまとめて購入できます',
            font_size=13, color=MUTED)

points = [
    '数量を選んで「カートに追加」',
    'カートアイコンで中身を確認',
    '数量の変更・削除もカート内で',
    '規格（バリエーション）がある場合は選択',
]
for i, pt in enumerate(points):
    y = Inches(1.8 + i * 0.55)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Emu(y.emu + Inches(0.05).emu),
                              Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = MAIN
    dot.line.fill.background()
    add_textbox(s, Inches(1.3), y, Inches(4.5), Inches(0.4),
                pt, font_size=13, color=TEXT_CLR)

info_box = add_rounded_rect(s, Inches(0.8), Inches(4.3), Inches(5.4), Inches(0.7), HIGHLIGHT)
add_textbox(s, Inches(1.0), Inches(4.4), Inches(5.0), Inches(0.5),
            '複数の出品者の商品もまとめてカートに入れられます',
            font_size=12, color=MAIN_DARK, bold=True, alignment=PP_ALIGN.LEFT)

add_screenshot_placeholder(s, Inches(6.8), Inches(1.7), Inches(5.8), Inches(5.2),
                           'カート画面')
add_stripe_bottom(s)
add_slide_num(s, 11, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — STEP 3: 購入手続き
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
step_label = add_rounded_rect(s, Inches(0.8), Inches(0.5), Inches(1.2), Inches(0.4), MAIN)
tf = step_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = 'STEP 3'
r.font.size = Pt(11)
r.font.color.rgb = WHITE
r.font.bold = True
r.font.name = FONT_JP

add_textbox(s, Inches(2.2), Inches(0.45), Inches(6), Inches(0.5),
            '購入手続き', font_size=26, bold=True, color=MAIN_DARK)
add_textbox(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.3),
            '受け取り方法・配送先・お届け日時・支払い方法を選択',
            font_size=13, color=MUTED)

add_screenshot_placeholder(s, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.2),
                           'チェックアウト画面')

steps = [
    ('1', '受け取り方法を選ぶ', '配送 or 畑受け取り'),
    ('2', '配送先を入力 or 選択', '以前の住所は自動表示'),
    ('3', 'お届け日時を選ぶ', '希望の日時を指定'),
    ('4', '支払い方法を選ぶ', 'カード / 掛売 / 代引き'),
    ('5', '注文を確定する', '内容を確認して完了'),
]
for i, (n, t, d) in enumerate(steps):
    y = Inches(1.9 + i * 0.95)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), y, Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = MAIN
    circ.line.fill.background()
    tf_c = circ.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    r_c = p_c.add_run()
    r_c.text = n
    r_c.font.size = Pt(13)
    r_c.font.color.rgb = WHITE
    r_c.font.bold = True

    add_textbox(s, Inches(7.8), y, Inches(4.5), Inches(0.3),
                t, font_size=14, bold=True, color=MAIN_DARK)
    add_textbox(s, Inches(7.8), Emu(y.emu + Inches(0.3).emu), Inches(4.5), Inches(0.25),
                d, font_size=11, color=MUTED)

add_stripe_bottom(s)
add_slide_num(s, 12, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — STEP 4: 注文確定
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
step_label = add_rounded_rect(s, Inches(0.8), Inches(0.5), Inches(1.2), Inches(0.4), MAIN)
tf = step_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = 'STEP 4'
r.font.size = Pt(11)
r.font.color.rgb = WHITE
r.font.bold = True
r.font.name = FONT_JP

add_textbox(s, Inches(2.2), Inches(0.45), Inches(6), Inches(0.5),
            '注文確定', font_size=26, bold=True, color=MAIN_DARK)
add_textbox(s, Inches(0.8), Inches(1.1), Inches(8), Inches(0.3),
            '注文完了画面が表示され、メールで注文確認が届きます',
            font_size=13, color=MUTED)

add_screenshot_placeholder(s, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.2),
                           '注文確認画面')
add_screenshot_placeholder(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.2),
                           '注文完了画面')

add_stripe_bottom(s)
add_slide_num(s, 13, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — 選べる決済方法
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
add_textbox(s, Inches(1), Inches(0.4), Inches(11.3), Inches(0.6),
            '選べる決済方法', font_size=32, bold=True, color=MAIN_DARK,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER)

s.shapes.add_picture(IMG('slide_payment.png'),
                     Inches(0.8), Inches(1.3), Inches(11.7), Inches(3.5))

payments = [
    ('クレジットカード', '注文時にカード情報を入力。\nStripe社の安全な決済基盤。', '個人・法人 どちらも利用可', MAIN),
    ('掛売（請求書払い）', '月末締め・翌月末払い。\n飲食店の実務に最適な決済。', '法人のお客様向け', ACCENT),
    ('代金引換', '商品のお届け時に\n配送員へお支払い。', '個人・法人 どちらも利用可', MAIN),
]
for i, (t, d, badge_text, c) in enumerate(payments):
    x = Inches(0.8 + i * 4.1)
    add_textbox(s, x, Inches(4.95), Inches(3.5), Inches(0.35),
                t, font_size=16, bold=True, color=MAIN_DARK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, Inches(5.35), Inches(3.5), Inches(0.7),
                d, font_size=11, color=MUTED,
                alignment=PP_ALIGN.CENTER, line_spacing=1.5)
    badge = add_rounded_rect(s, Emu(x.emu + Inches(0.6).emu), Inches(6.3),
                              Inches(2.3), Inches(0.35), c)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = badge_text
    r.font.size = Pt(10)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = FONT_JP

add_stripe_bottom(s)
add_slide_num(s, 14, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — 掛売（請求書払い）の流れ
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
step_label = add_rounded_rect(s, Inches(0.8), Inches(0.5), Inches(1.6), Inches(0.4), ACCENT)
tf = step_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = '法人向け'
r.font.size = Pt(11)
r.font.color.rgb = MAIN_DARK
r.font.bold = True
r.font.name = FONT_JP

add_textbox(s, Inches(2.6), Inches(0.45), Inches(8), Inches(0.5),
            '掛売（請求書払い）の流れ', font_size=26, bold=True, color=MAIN_DARK)

invoice_steps = [
    ('1', '支払方法で「掛売」を選択して注文'),
    ('2', '月中に複数回の注文もOK'),
    ('3', '月末締めで請求書が発行される'),
    ('4', '請求書一覧からPDF確認・ダウンロード'),
    ('5', '翌月末までに銀行振込でお支払い'),
]
for i, (n, t) in enumerate(invoice_steps):
    y = Inches(1.5 + i * 1.0)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = MAIN
    circ.line.fill.background()
    tf_c = circ.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    r_c = p_c.add_run()
    r_c.text = n
    r_c.font.size = Pt(16)
    r_c.font.color.rgb = WHITE
    r_c.font.bold = True

    add_textbox(s, Inches(1.7), Emu(y.emu + Inches(0.08).emu), Inches(4.5), Inches(0.4),
                t, font_size=15, color=TEXT_CLR)

    if i < len(invoice_steps) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                    Inches(1.15), Emu(y.emu + Inches(0.55).emu),
                                    Inches(0.2), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

add_screenshot_placeholder(s, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.6),
                           '請求書一覧画面')
add_stripe_bottom(s)
add_slide_num(s, 15, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — 配送状況の確認
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
step_label = add_rounded_rect(s, Inches(0.8), Inches(0.5), Inches(1.6), Inches(0.4), MAIN)
tf = step_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = '注文後'
r.font.size = Pt(11)
r.font.color.rgb = WHITE
r.font.bold = True
r.font.name = FONT_JP

add_textbox(s, Inches(2.6), Inches(0.45), Inches(8), Inches(0.5),
            '配送状況の確認', font_size=26, bold=True, color=MAIN_DARK)
add_textbox(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.3),
            'ダッシュボードから注文のステータスをリアルタイムで確認できます',
            font_size=13, color=MUTED)

add_screenshot_placeholder(s, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.2),
                           '注文詳細 / 配送ステータス画面')

statuses = [
    ('出荷準備中', '商品を準備しています'),
    ('お届け中', '配送業者に渡しました'),
    ('配達完了', 'お届けが完了しました'),
]
header_bg = add_rounded_rect(s, Inches(7.2), Inches(2.0), Inches(5.2), Inches(0.5), MAIN)
tf = header_bg.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = '配送ステータスの見方'
r.font.size = Pt(12)
r.font.color.rgb = WHITE
r.font.bold = True
r.font.name = FONT_JP

for i, (label, desc) in enumerate(statuses):
    y = Inches(2.7 + i * 0.7)
    row_bg = add_rounded_rect(s, Inches(7.2), y, Inches(5.2), Inches(0.55),
                               SOFT if i % 2 == 0 else WHITE)
    add_textbox(s, Inches(7.4), y, Inches(1.8), Inches(0.5),
                label, font_size=13, bold=True, color=MAIN_DARK)
    add_textbox(s, Inches(9.3), y, Inches(3.0), Inches(0.5),
                desc, font_size=12, color=MUTED)

info_box = add_rounded_rect(s, Inches(7.2), Inches(5.0), Inches(5.2), Inches(0.8), HIGHLIGHT)
add_textbox(s, Inches(7.4), Inches(5.1), Inches(4.8), Inches(0.6),
            '注文が入るとメール通知が届き、\nダッシュボードからもいつでも確認できます。',
            font_size=12, color=MAIN_DARK, bold=True, line_spacing=1.5)

add_stripe_bottom(s)
add_slide_num(s, 16, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — 畑での受け取り
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
step_label = add_rounded_rect(s, Inches(0.8), Inches(0.5), Inches(3.0), Inches(0.4), ACCENT)
tf = step_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = 'セッツマルシェならでは'
r.font.size = Pt(11)
r.font.color.rgb = MAIN_DARK
r.font.bold = True
r.font.name = FONT_JP

add_textbox(s, Inches(4.0), Inches(0.45), Inches(6), Inches(0.5),
            '畑での受け取り', font_size=26, bold=True, color=MAIN_DARK)

s.shapes.add_picture(IMG('slide_farm_pickup.png'),
                     Inches(0.8), Inches(1.3), Inches(6.0), Inches(5.5))

pickup_steps = [
    ('1', 'チェックアウトで「畑で受け取り」を選択', '配送先の入力は不要です'),
    ('2', '出品者から準備完了の連絡', '受け取り日時を調整します'),
    ('3', '畑で直接受け取り', '生産者と直接会える体験'),
]
for i, (n, t, d) in enumerate(pickup_steps):
    y = Inches(1.5 + i * 1.3)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), y, Inches(0.45), Inches(0.45))
    circ.fill.solid()
    circ.fill.fore_color.rgb = MAIN
    circ.line.fill.background()
    tf_c = circ.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    r_c = p_c.add_run()
    r_c.text = n
    r_c.font.size = Pt(15)
    r_c.font.color.rgb = WHITE
    r_c.font.bold = True

    add_textbox(s, Inches(7.9), y, Inches(4.5), Inches(0.35),
                t, font_size=14, bold=True, color=MAIN_DARK)
    add_textbox(s, Inches(7.9), Emu(y.emu + Inches(0.35).emu), Inches(4.5), Inches(0.3),
                d, font_size=11, color=MUTED)

merits = add_rounded_rect(s, Inches(7.0), Inches(5.5), Inches(5.8), Inches(1.0), HIGHLIGHT)
add_textbox(s, Inches(7.2), Inches(5.6), Inches(5.4), Inches(0.8),
            '生産者の顔が見える安心感と、\n配送コストゼロのお得さ。',
            font_size=14, color=MAIN_DARK, bold=True,
            alignment=PP_ALIGN.CENTER, line_spacing=1.6)

add_stripe_bottom(s)
add_slide_num(s, 17, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — 便利な機能
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
add_textbox(s, Inches(1), Inches(0.4), Inches(11.3), Inches(0.6),
            '便利な機能', font_size=32, bold=True, color=MAIN_DARK,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(2.5), Inches(1.05), Inches(8.3), Inches(0.4),
            '毎日の仕入れをもっとラクに', font_size=15, color=MUTED,
            alignment=PP_ALIGN.CENTER)

s.shapes.add_picture(IMG('slide_convenience.png'),
                     Inches(1.0), Inches(1.6), Inches(11.3), Inches(4.0))

features = [
    ('リピート注文', '過去の注文から「この注文をもう一度」\nボタンひとつでカートに追加。', '注文履歴から1クリック'),
    ('注文テンプレート', 'よく注文する商品の組み合わせを保存。\n「いつもの」をワンクリックで注文。', 'カートから保存も可能'),
]
for i, (t, d, badge) in enumerate(features):
    x = Inches(1.2 + i * 6.0)
    add_textbox(s, x, Inches(5.75), Inches(5.0), Inches(0.35),
                t, font_size=17, bold=True, color=MAIN_DARK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, Inches(6.15), Inches(5.0), Inches(0.7),
                d, font_size=11, color=MUTED,
                alignment=PP_ALIGN.CENTER, line_spacing=1.5)

add_stripe_bottom(s)
add_slide_num(s, 18, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 — 法人向け安心機能
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_bg_color(s, WHITE)
add_textbox(s, Inches(1), Inches(0.4), Inches(11.3), Inches(0.6),
            '法人向けの安心機能', font_size=32, bold=True, color=MAIN_DARK,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(2.5), Inches(1.05), Inches(8.3), Inches(0.4),
            '組織での利用に対応した仕組みをご用意しています',
            font_size=15, color=MUTED, alignment=PP_ALIGN.CENTER)

biz_features = [
    ('顧客別価格', 'お取引状況に応じた特別価格を\n設定。ログインすると自動適用。', MAIN),
    ('購買承認フロー', '担当者が注文 → 店長が承認。\n組織に合わせたステップを設定。', ACCENT),
    ('組織メンバー管理', '注文担当・承認者・経理など、\nメンバーごとに権限を設定。', HIGHLIGHT),
]
for i, (t, d, c) in enumerate(biz_features):
    x = Inches(0.8 + i * 4.1)
    card_bg = add_rounded_rect(s, x, Inches(1.8), Inches(3.7), Inches(4.8),
                                WHITE, RGBColor(0xE5,0xE7,0xEB), Pt(0.75))
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(3.7), Pt(5))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = c
    top_bar.line.fill.background()

    icon_area = add_rounded_rect(s, Emu(x.emu + Inches(1.2).emu), Inches(2.3),
                                  Inches(1.3), Inches(1.3), SOFT)
    icons = ['¥', '✓', '⚙']
    tf_i = icon_area.text_frame
    p_i = tf_i.paragraphs[0]
    p_i.alignment = PP_ALIGN.CENTER
    r_i = p_i.add_run()
    r_i.text = icons[i]
    r_i.font.size = Pt(36)
    r_i.font.color.rgb = MAIN
    r_i.font.bold = True

    add_textbox(s, Emu(x.emu + Inches(0.3).emu), Inches(3.9),
                Inches(3.1), Inches(0.4),
                t, font_size=18, bold=True, color=MAIN_DARK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(s, Emu(x.emu + Inches(0.3).emu), Inches(4.5),
                Inches(3.1), Inches(0.8),
                d, font_size=12, color=MUTED,
                alignment=PP_ALIGN.CENTER, line_spacing=1.6)

add_stripe_bottom(s)
add_slide_num(s, 19, TOTAL)

# ═══════════════════════════════════════════════════════════════
# SLIDE 20 — クロージング
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_image_bg(s, IMG('slide_vision.png'))
add_overlay(s, MAIN_DARK, alpha_pct=60)

add_textbox(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(2.0),
            '地域でつくり、地域で食べる。\nその循環を、ごいっしょに。',
            font_size=36, bold=True, color=WHITE,
            font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

add_textbox(s, Inches(2.5), Inches(3.8), Inches(8.3), Inches(1.0),
            'セッツマルシェは、飲食店のみなさまと一緒に育てていくサービスです。\nご意見・ご要望をぜひお聞かせください。',
            font_size=15, color=RGBColor(0xDD,0xDD,0xDD),
            alignment=PP_ALIGN.CENTER, line_spacing=1.7)

contact_box = add_rounded_rect(s, Inches(3.5), Inches(5.3), Inches(6.3), Inches(1.3),
                                RGBColor(0xFF,0xFF,0xFF))
from lxml import etree
sp_elem = contact_box._element
solid = sp_elem.find('.//' + qn('a:solidFill'))
if solid is not None:
    srgb = solid.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha_el = etree.SubElement(srgb, qn('a:alpha'))
        alpha_el.set('val', str(15 * 1000))

add_textbox(s, Inches(3.8), Inches(5.5), Inches(5.7), Inches(0.4),
            'セッツマルシェ（渡辺ファーム）', font_size=15, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.8), Inches(6.0), Inches(5.7), Inches(0.4),
            'お問い合わせ：サイト内「お問い合わせフォーム」より',
            font_size=12, color=RGBColor(0xCC,0xCC,0xCC),
            alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
out = os.path.join(DIR, 'セッツマルシェ_ご紹介資料.pptx')
prs.save(out)
print(f'✅ Saved: {out}')
print(f'   Slides: {len(prs.slides)}')
