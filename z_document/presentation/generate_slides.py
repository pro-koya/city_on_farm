#!/usr/bin/env python3
"""セッツマルシェ ご紹介資料 — pptx generator v3
   Story-driven, infographic-heavy, minimal AI illustrations."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from lxml import etree

BASE = os.path.dirname(os.path.abspath(__file__))
OUT  = os.path.join(BASE, "セッツマルシェ_ご紹介資料.pptx")

# ─── Brand tokens ──────────────────────────────────────────────
MAIN      = RGBColor(0x4C, 0x6B, 0x5C)
MAIN_DARK = RGBColor(0x35, 0x5D, 0x4A)
ACCENT    = RGBColor(0xA3, 0xC9, 0xA8)
HIGHLIGHT = RGBColor(0xD6, 0xEA, 0xDF)
SOFT      = RGBColor(0xF2, 0xF7, 0xF4)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_CLR  = RGBColor(0x2B, 0x2B, 0x2B)
MUTED     = RGBColor(0x6B, 0x72, 0x80)
DARK_BG   = RGBColor(0x2D, 0x42, 0x36)
WARM      = RGBColor(0xE8, 0x92, 0x7C)
WARM_SOFT = RGBColor(0xFD, 0xF0, 0xEB)
WARM_BG   = RGBColor(0xFB, 0xF5, 0xF1)
GOLD      = RGBColor(0xC4, 0x8E, 0x5A)
BLUE_ACC  = RGBColor(0x6B, 0xA3, 0xD6)
BORDER    = RGBColor(0xE5, 0xE7, 0xEB)

FONT   = "Yu Gothic"
SERIF  = "Yu Mincho"
NS     = 'http://schemas.openxmlformats.org/drawingml/2006/main'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

img_path = lambda n: os.path.join(BASE, n)
slide_n = [0]

# ═══════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════

def new_slide():
    slide_n[0] += 1
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(sl, c):
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = c

def rect(sl, l, t, w, h, c, alpha=None):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    if alpha:
        sf = s._element.find(f'.//{{{NS}}}solidFill')
        if sf is not None:
            etree.SubElement(sf[0], f'{{{NS}}}alpha').set('val', str(int(alpha*1000)))
    return s

def rrect(sl, l, t, w, h, c, border_c=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    if border_c:
        s.line.color.rgb = border_c; s.line.width = Pt(1)
        s.line.fill.solid(); s.line.fill.fore_color.rgb = border_c
    return s

def circle(sl, l, t, sz, c):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s

def _vcenter(shp):
    bp = shp._element.find(f'.//{{{NS}}}bodyPr')
    if bp is not None: bp.set('anchor', 'ctr')

def tb(sl, l, t, w, h, txt, sz=18, c=TEXT_CLR, b=False, a=PP_ALIGN.LEFT,
       f=FONT, ls=None, va=False):
    box = sl.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = f; p.alignment = a
    p.space_before = p.space_after = Pt(0)
    if ls: p.line_spacing = Pt(ls)
    if va: _vcenter(box)
    return box

def mtb(sl, l, t, w, h, lines, a=PP_ALIGN.LEFT, f=FONT):
    """lines: list of (text, size, color, bold)"""
    box = sl.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, sz, clr, bld) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = clr
        p.font.bold = bld; p.font.name = f; p.alignment = a
        p.space_before = Pt(4); p.space_after = Pt(4)
    return box

def pill(sl, l, t, w, h, txt, bg_c=MAIN, fg=WHITE, sz=12, b=True):
    s = rrect(sl, l, t, w, h, bg_c)
    tf = s.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = fg; p.font.bold = b
    p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    _vcenter(s)
    return s

def clabel(sl, l, t, sz_i, bg_c, txt, tsz=15, tc=WHITE):
    s = circle(sl, l, t, sz_i, bg_c)
    tf = s.text_frame; tf.word_wrap = False; tf.auto_size = None
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(tsz); p.font.color.rgb = tc; p.font.bold = True
    p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    _vcenter(s)
    return s

def stripe(sl):
    bh = Inches(0.06); y = SH - bh; th = SW // 3
    rect(sl, Emu(0), y, th, bh, MAIN)
    rect(sl, th, y, th, bh, ACCENT)
    rect(sl, th*2, y, th + Emu(200), bh, HIGHLIGHT)

def snum(sl):
    tb(sl, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
       f"{slide_n[0]}", sz=10, c=MUTED, a=PP_ALIGN.RIGHT)

def img_fill(sl, path):
    if os.path.exists(path):
        sl.shapes.add_picture(path, Emu(0), Emu(0), SW, SH)

def img_prop(sl, path, l, t, width=None, height=None):
    if not os.path.exists(path): return
    im = Image.open(path); ratio = im.size[0] / im.size[1]
    if width and not height:  h = int(width / ratio);  sl.shapes.add_picture(path, l, t, width, h)
    elif height and not width: w = int(height * ratio); sl.shapes.add_picture(path, l, t, w, height)

def ss_placeholder(sl, l, t, w, h, label):
    s = rrect(sl, l, t, w, h, SOFT)
    s.line.color.rgb = ACCENT; s.line.width = Pt(1.5); s.line.dash_style = 4
    s.line.fill.solid(); s.line.fill.fore_color.rgb = ACCENT
    tb(sl, l, t + h//2 - Inches(0.3), w, Inches(0.3), label, sz=14, c=MUTED, a=PP_ALIGN.CENTER)
    tb(sl, l, t + h//2, w, Inches(0.25),
       "[ スクリーンショットを挿入 ]", sz=11, c=ACCENT, a=PP_ALIGN.CENTER)

def arrow_right(sl, l, t, w, c=ACCENT):
    s = sl.shapes.add_shape(MSO_SHAPE.NOTCHED_RIGHT_ARROW, l, t, w, Inches(0.35))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s

def heading(sl, txt, sub=None, y_title=Inches(0.35)):
    tb(sl, Inches(0.8), y_title, Inches(11.5), Inches(0.6),
       txt, sz=28, c=MAIN_DARK, b=True, a=PP_ALIGN.CENTER, f=SERIF)
    if sub:
        tb(sl, Inches(1.5), y_title + Inches(0.6), Inches(10), Inches(0.4),
           sub, sz=15, c=MUTED, a=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
img_fill(sl, img_path("slide_hero.png"))
rect(sl, Emu(0), Emu(0), SW, SH, DARK_BG, alpha=55)

tb(sl, Inches(1.5), Inches(1.2), Inches(10), Inches(0.5),
   "2026.03.24  ご紹介資料", sz=14, c=WHITE, a=PP_ALIGN.CENTER)
tb(sl, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
   "セッツマルシェ", sz=56, c=WHITE, a=PP_ALIGN.CENTER, b=True, f=SERIF)
tb(sl, Inches(1.5), Inches(3.4), Inches(10), Inches(0.8),
   "まちと畑を、ひとつの食卓につなぐ。", sz=24, c=WHITE, a=PP_ALIGN.CENTER, f=SERIF)

rect(sl, Inches(5.5), Inches(4.6), Inches(2.3), Inches(0.03), WHITE)

tb(sl, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
   "渡辺ファーム  ／  摂津市商工会", sz=15, c=WHITE, a=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

heading(sl, "本日の内容")

items = [
    ("1", "課題と背景", "飲食店・農家が抱える「いま」"),
    ("2", "セッツマルシェの提案", "コンセプト・3つの柱・目指す姿"),
    ("3", "実際の購入フロー", "購入手順・決済・配送・畑受け取り"),
    ("4", "便利な機能と今後", "リピート注文・テンプレート・法人機能"),
]
y = Inches(1.7)
for num, title, sub in items:
    clabel(sl, Inches(4.0), y, Inches(0.5), MAIN, num, tsz=16)
    tb(sl, Inches(4.8), y + Inches(0.02), Inches(5), Inches(0.3),
       title, sz=19, c=MAIN_DARK, b=True)
    tb(sl, Inches(4.8), y + Inches(0.36), Inches(5), Inches(0.25),
       sub, sz=13, c=MUTED)
    if num != "4":
        rect(sl, Inches(4.24), y + Inches(0.55), Inches(0.02), Inches(0.55), ACCENT)
    y += Inches(1.1)
snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 3 — HOOK: "こんな経験、ありませんか？"
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WARM_BG)

tb(sl, Inches(1), Inches(2.4), Inches(11), Inches(1.2),
   "こんな経験、ありませんか？", sz=40, c=WARM, b=True, a=PP_ALIGN.CENTER, f=SERIF)

tb(sl, Inches(2), Inches(4.0), Inches(9), Inches(0.6),
   "ー  飲食店を営むあなたへ  ー", sz=16, c=MUTED, a=PP_ALIGN.CENTER)

stripe(sl); snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 4 — PROBLEM 1: 飲食店の声
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

rect(sl, Emu(0), Emu(0), Inches(0.12), SH, WARM)

tb(sl, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4),
   "飲食店の現実", sz=13, c=WARM, b=True)

tb(sl, Inches(0.8), Inches(1.2), Inches(6), Inches(1.5),
   "「新鮮な地場野菜が\n  手に入らない」", sz=34, c=TEXT_CLR, b=True, f=SERIF, ls=52)

tb(sl, Inches(0.8), Inches(3.2), Inches(5.5), Inches(1.2),
   "大手卸では産地が遠い。\n地元農家から直接買いたくても、\n接点がなく、安定供給の保証もない。",
   sz=15, c=MUTED, ls=28)

# Right side: infographic - pain points
card_y = Inches(0.8)
pain = [
    ("仕入れ先が限定的", "大手卸・市場に頼るしかなく\n地場野菜の選択肢がない"),
    ("鮮度のタイムラグ", "収穫から店頭まで時間がかかり\n「採れたて」が届かない"),
    ("コストの壁", "少量ロットの仕入れでは\n配送コストが割高になる"),
]
for title, desc in pain:
    card = rrect(sl, Inches(7.2), card_y, Inches(5.3), Inches(1.7), WHITE, border_c=BORDER)
    rect(sl, Inches(7.2), card_y, Inches(0.07), Inches(1.7), WARM)
    tb(sl, Inches(7.5), card_y + Inches(0.2), Inches(4.8), Inches(0.35),
       title, sz=15, c=WARM, b=True)
    tb(sl, Inches(7.5), card_y + Inches(0.65), Inches(4.8), Inches(0.7),
       desc, sz=12, c=MUTED)
    card_y += Inches(1.95)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 5 — PROBLEM 2: 農家の声
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

rect(sl, Emu(0), Emu(0), Inches(0.12), SH, WARM)

tb(sl, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4),
   "農家の現実", sz=13, c=WARM, b=True)

tb(sl, Inches(0.8), Inches(1.2), Inches(6), Inches(1.5),
   "「せっかく育てたのに\n  届ける先がない」", sz=34, c=TEXT_CLR, b=True, f=SERIF, ls=52)

tb(sl, Inches(0.8), Inches(3.2), Inches(5.5), Inches(1.2),
   "小規模農家や市民農園では、\n少量の余剰を販売するルートがなく、\nせっかくの収穫が廃棄されてしまう。",
   sz=15, c=MUTED, ls=28)

# Right side: numbers impact
data_cards = [
    ("販路なし", "小規模農家の多くが\n直売以外の販路を持たない"),
    ("余剰の廃棄", "市民農園の収穫物は\n自家消費しきれず廃棄に"),
    ("つながりの不在", "飲食店と農家の間に\n直接取引の接点がない"),
]
cy = Inches(0.8)
for title, desc in data_cards:
    card = rrect(sl, Inches(7.2), cy, Inches(5.3), Inches(1.7), WHITE, border_c=BORDER)
    rect(sl, Inches(7.2), cy, Inches(0.07), Inches(1.7), WARM)
    tb(sl, Inches(7.5), cy + Inches(0.2), Inches(4.8), Inches(0.35),
       title, sz=15, c=WARM, b=True)
    tb(sl, Inches(7.5), cy + Inches(0.65), Inches(4.8), Inches(0.7),
       desc, sz=12, c=MUTED)
    cy += Inches(1.95)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 6 — BRIDGE: "もし、こんな仕組みがあったら？"
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, SOFT)

tb(sl, Inches(1), Inches(1.8), Inches(11), Inches(1),
   "もし、こんな仕組みがあったら？", sz=36, c=MAIN_DARK, b=True, a=PP_ALIGN.CENTER, f=SERIF)

hopes = [
    "地元の新鮮な野菜を、スマホひとつで注文できる",
    "小さな農家の野菜も、飲食店の一皿として輝く",
    "畑で直接受け取れば、採れたての美味しさをそのまま",
    "生産者の顔が見える安心感",
]
y = Inches(3.4)
for h in hopes:
    clabel(sl, Inches(3.5), y + Inches(0.02), Inches(0.32), ACCENT, "✓", tsz=14, tc=WHITE)
    tb(sl, Inches(4.1), y, Inches(6), Inches(0.4),
       h, sz=17, c=MAIN_DARK)
    y += Inches(0.6)

stripe(sl); snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 7 — ANSWER: "それが、セッツマルシェです"
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, MAIN_DARK)

tb(sl, Inches(1), Inches(2.6), Inches(11), Inches(1.2),
   "それが、セッツマルシェです。", sz=42, c=WHITE, b=True, a=PP_ALIGN.CENTER, f=SERIF)

tb(sl, Inches(2), Inches(4.2), Inches(9), Inches(0.6),
   "まちと畑を、ひとつの食卓につなぐオンラインマルシェ", sz=17, c=ACCENT, a=PP_ALIGN.CENTER)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 8 — HOW IT WORKS: diagram (playful design)
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

# Decorative background shapes for visual interest
circle(sl, Inches(-0.5), Inches(-0.5), Inches(3), SOFT)
circle(sl, Inches(11), Inches(5.5), Inches(3), SOFT)
rect(sl, Inches(4.2), Inches(0.8), Inches(4.9), Inches(6.5), HIGHLIGHT)

heading(sl, "セッツマルシェの仕組み")

# Left column label
pill(sl, Inches(0.5), Inches(1.5), Inches(2.7), Inches(0.4),
     "つくる人", bg_c=ACCENT, fg=MAIN_DARK, sz=13)

prod_y = Inches(2.2)
producers = [
    (MAIN, "渡辺ファーム", "主力野菜を安定供給"),
    (ACCENT, "小規模農家", "少量でも出品OK"),
    (GOLD, "市民農園", "余剰を活用"),
]
for clr, name, desc in producers:
    c = rrect(sl, Inches(0.3), prod_y, Inches(3.1), Inches(1.2), WHITE, border_c=BORDER)
    clabel(sl, Inches(0.5), prod_y + Inches(0.2), Inches(0.7), clr, name[0], tsz=18)
    tb(sl, Inches(1.4), prod_y + Inches(0.18), Inches(1.8), Inches(0.3),
       name, sz=14, c=MAIN_DARK, b=True)
    tb(sl, Inches(1.4), prod_y + Inches(0.55), Inches(1.8), Inches(0.3),
       desc, sz=11, c=MUTED)
    prod_y += Inches(1.4)

# Center: セッツマルシェ hub (larger, more prominent)
arrow_right(sl, Inches(3.5), Inches(3.5), Inches(0.9), MAIN)

hub = rrect(sl, Inches(4.5), Inches(1.8), Inches(4.3), Inches(4.3), MAIN)
circle(sl, Inches(5.75), Inches(2.2), Inches(1.8), RGBColor(0x3A, 0x55, 0x45))
tb(sl, Inches(5.75), Inches(2.6), Inches(1.8), Inches(1.0),
   "セッツ\nマルシェ", sz=16, c=WHITE, b=True, a=PP_ALIGN.CENTER, f=SERIF)

hub_features = [
    ("商品掲載", Inches(4.8), Inches(4.3)),
    ("受注管理", Inches(6.6), Inches(4.3)),
    ("決済処理", Inches(4.8), Inches(4.85)),
    ("配送手配", Inches(6.6), Inches(4.85)),
]
for feat, fx, fy in hub_features:
    pill(sl, fx, fy, Inches(1.6), Inches(0.4),
         feat, bg_c=RGBColor(0x3A, 0x55, 0x45), fg=ACCENT, sz=11, b=False)

tb(sl, Inches(4.5), Inches(5.5), Inches(4.3), Inches(0.4),
   "オンラインでつなぐ仕組み", sz=12, c=ACCENT, a=PP_ALIGN.CENTER)

# Right column
arrow_right(sl, Inches(8.9), Inches(3.5), Inches(0.9), MAIN)

pill(sl, Inches(10.1), Inches(1.5), Inches(2.7), Inches(0.4),
     "届く人", bg_c=GOLD, fg=WHITE, sz=13)

buyer_y = Inches(2.2)
buyers = [
    (GOLD, "飲食店", "法人のお客様", "定期仕入れ"),
    (BLUE_ACC, "ご家庭", "個人のお客様", "旬の野菜を"),
]
for clr, name, desc, note in buyers:
    c = rrect(sl, Inches(9.9), buyer_y, Inches(3.1), Inches(1.8), WHITE, border_c=BORDER)
    clabel(sl, Inches(10.1), buyer_y + Inches(0.35), Inches(0.7), clr, name[0], tsz=18)
    tb(sl, Inches(11.0), buyer_y + Inches(0.25), Inches(1.8), Inches(0.3),
       name, sz=15, c=MAIN_DARK, b=True)
    tb(sl, Inches(11.0), buyer_y + Inches(0.6), Inches(1.8), Inches(0.25),
       desc, sz=11, c=MUTED)
    tb(sl, Inches(11.0), buyer_y + Inches(0.9), Inches(1.8), Inches(0.25),
       note, sz=11, c=clr, b=True)
    buyer_y += Inches(2.1)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 9 — PILLAR 1: 小さな声を拾う
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

pill(sl, Inches(0.8), Inches(0.4), Inches(1.2), Inches(0.35),
     "柱  1", bg_c=MAIN, sz=12)
tb(sl, Inches(2.2), Inches(0.33), Inches(6), Inches(0.5),
   "小さな声を拾う", sz=28, c=MAIN_DARK, b=True, f=SERIF)

# Before / After comparison
tb(sl, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
   "BEFORE — いま", sz=14, c=WARM, b=True)
rect(sl, Inches(0.8), Inches(1.75), Inches(5.5), Inches(0.03), WARM)

before_items = [
    "少量の収穫 → 販路がなく自家消費 or 廃棄",
    "市民農園の野菜 → 近所に配るだけ",
    "規格外品 → 商品にならないと判断",
]
by = Inches(2.0)
for item in before_items:
    tb(sl, Inches(1.0), by, Inches(5.3), Inches(0.35),
       "×  " + item, sz=13, c=MUTED)
    by += Inches(0.45)

tb(sl, Inches(0.8), Inches(3.7), Inches(5.5), Inches(0.4),
   "AFTER — セッツマルシェなら", sz=14, c=MAIN, b=True)
rect(sl, Inches(0.8), Inches(4.15), Inches(5.5), Inches(0.03), MAIN)

after_items = [
    "少量でもサイトに出品 → 飲食店が購入",
    "市民農園の余剰 → プロの料理で「一皿」に",
    "規格外品 → 「味は同じ」と理解ある買い手へ",
]
ay = Inches(4.4)
for item in after_items:
    tb(sl, Inches(1.0), ay, Inches(5.3), Inches(0.35),
       "✓  " + item, sz=13, c=MAIN_DARK)
    ay += Inches(0.45)

# Right side: visual summary
rect(sl, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.8), SOFT)
tb(sl, Inches(7.2), Inches(1.6), Inches(5.0), Inches(0.5),
   "対象となる生産者", sz=16, c=MAIN_DARK, b=True, a=PP_ALIGN.CENTER)

src_data = [
    (MAIN, "渡辺ファーム", "安定した主力野菜の供給源\n運営元として品質を保証"),
    (ACCENT, "小規模農家", "地域で丁寧に育てた野菜\n少量でも出品可能な仕組み"),
    (GOLD, "市民農園", "家庭菜園の余剰を活用\n小さな生産にも価値を"),
]
sy = Inches(2.3)
for clr, name, desc in src_data:
    rrect(sl, Inches(7.4), sy, Inches(4.6), Inches(1.35), WHITE, border_c=BORDER)
    clabel(sl, Inches(7.6), sy + Inches(0.25), Inches(0.7), clr, name[0], tsz=18)
    tb(sl, Inches(8.5), sy + Inches(0.15), Inches(3.3), Inches(0.35),
       name, sz=14, c=MAIN_DARK, b=True)
    tb(sl, Inches(8.5), sy + Inches(0.55), Inches(3.3), Inches(0.6),
       desc, sz=11, c=MUTED)
    sy += Inches(1.55)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 10 — PILLAR 2: まちの食を支える
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

pill(sl, Inches(0.8), Inches(0.4), Inches(1.2), Inches(0.35),
     "柱  2", bg_c=MAIN, sz=12)
tb(sl, Inches(2.2), Inches(0.33), Inches(6), Inches(0.5),
   "まちの食を支える", sz=28, c=MAIN_DARK, b=True, f=SERIF)

# Diagram: delivery area concept
tb(sl, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
   "地域内限定配送で、新鮮さと安定を両立", sz=16, c=MUTED, a=PP_ALIGN.CENTER)

# Feature cards in 2x2 grid
features = [
    ("鮮度", "地域内配送に限定\n収穫から最短でお届け", MAIN),
    ("安定", "複数の生産者から集荷\n季節に合わせた品揃え", ACCENT),
    ("安心", "生産者の顔が見える\n栽培方法・産地が明確", GOLD),
    ("手軽", "スマホ・PCから注文\n24時間いつでも発注OK", BLUE_ACC),
]
positions = [
    (Inches(0.8), Inches(2.2)),
    (Inches(6.9), Inches(2.2)),
    (Inches(0.8), Inches(4.6)),
    (Inches(6.9), Inches(4.6)),
]
for (l, t), (title, desc, clr) in zip(positions, features):
    card = rrect(sl, l, t, Inches(5.6), Inches(2.0), WHITE, border_c=BORDER)
    rect(sl, l, t, Inches(5.6), Inches(0.06), clr)
    clabel(sl, l + Inches(0.3), t + Inches(0.3), Inches(0.65), clr,
           title[0], tsz=20, tc=WHITE)
    tb(sl, l + Inches(1.2), t + Inches(0.25), Inches(4), Inches(0.4),
       title, sz=18, c=MAIN_DARK, b=True)
    tb(sl, l + Inches(1.2), t + Inches(0.75), Inches(4), Inches(0.8),
       desc, sz=13, c=MUTED)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 11 — PILLAR 3: 本当の美味しさを届ける
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

pill(sl, Inches(0.8), Inches(0.4), Inches(1.2), Inches(0.35),
     "柱  3", bg_c=MAIN, sz=12)
tb(sl, Inches(2.2), Inches(0.33), Inches(8), Inches(0.5),
   "本当の美味しさを届ける", sz=28, c=MAIN_DARK, b=True, f=SERIF)

tb(sl, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5),
   "「美味しさ」は、味だけではありません。", sz=17, c=MUTED, a=PP_ALIGN.CENTER)

# Three layers of "oishisa"
layers = [
    (MAIN, "素材の美味しさ",
     "地元で採れた旬の野菜。\n収穫から最短でお届けすることで、\n鮮度そのものの味わいを届けます。"),
    (ACCENT, "つながりの美味しさ",
     "「この野菜は、あの農家さんが育てた」\nと語れる仕入れ。お客様への\nストーリーが料理の付加価値に。"),
    (GOLD, "循環の美味しさ",
     "地域で育て、地域で食べる。\nその循環を実感できることが、\n関わるすべての人の喜びになります。"),
]

x = Inches(0.5)
for clr, title, desc in layers:
    card = rrect(sl, x, Inches(2.0), Inches(3.9), Inches(4.8), WHITE, border_c=BORDER)
    rect(sl, x, Inches(2.0), Inches(3.9), Inches(0.07), clr)
    # big circle icon
    clabel(sl, x + Inches(1.35), Inches(2.5), Inches(1.2), clr,
           title[0], tsz=30, tc=WHITE)
    tb(sl, x + Inches(0.3), Inches(3.9), Inches(3.3), Inches(0.4),
       title, sz=17, c=MAIN_DARK, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.3), Inches(4.5), Inches(3.3), Inches(1.5),
       desc, sz=13, c=MUTED, a=PP_ALIGN.CENTER, ls=22)
    x += Inches(4.2)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 12 — VISION
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
img_fill(sl, img_path("slide_closing.png"))
rect(sl, Emu(0), Emu(0), SW, SH, DARK_BG, alpha=58)

tb(sl, Inches(1.5), Inches(2.0), Inches(10), Inches(1.8),
   "地域でつくり、\n地域で味わい、\n地域でつながる。",
   sz=40, c=WHITE, b=True, a=PP_ALIGN.CENTER, f=SERIF, ls=60)
tb(sl, Inches(2.5), Inches(4.6), Inches(8), Inches(0.8),
   "この循環の輪が、まちの豊かさを育んでいく。\nそれが、私たち「セッツマルシェ」の願いです。",
   sz=16, c=WHITE, a=PP_ALIGN.CENTER, ls=28)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 13 — PART 2 DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, MAIN_DARK)

pill(sl, Inches(5.6), Inches(2.3), Inches(2.1), Inches(0.45),
     "PART 2", bg_c=ACCENT, fg=MAIN_DARK, sz=14)

tb(sl, Inches(1.5), Inches(3.1), Inches(10), Inches(0.8),
   "ご利用方法のご紹介", sz=38, c=WHITE, b=True, a=PP_ALIGN.CENTER, f=SERIF)
tb(sl, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
   "飲食店さま向け — 購入フローと便利な機能", sz=16, c=ACCENT, a=PP_ALIGN.CENTER)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 14 — PURCHASE FLOW OVERVIEW (playful diagram)
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, SOFT); stripe(sl)

# Decorative background curves
circle(sl, Inches(10.5), Inches(-1), Inches(4), HIGHLIGHT)
circle(sl, Inches(-1.5), Inches(5), Inches(4), HIGHLIGHT)

heading(sl, "購入の流れ — 全体像")

# Step flow with varied colors and icons
flow_steps = [
    ("商品を\n探す", MAIN, "1"),
    ("カートに\n入れる", RGBColor(0x5A, 0x8A, 0x70), "2"),
    ("購入\n手続き", ACCENT, "3"),
    ("決済", GOLD, "4"),
    ("注文\n確定", MAIN_DARK, "5"),
]
x = Inches(0.5)
for label, clr, num in flow_steps:
    # Card with top circle
    card = rrect(sl, x, Inches(2.2), Inches(2.1), Inches(2.4), WHITE, border_c=BORDER)
    rect(sl, x, Inches(2.2), Inches(2.1), Inches(0.07), clr)
    clabel(sl, x + Inches(0.65), Inches(1.75), Inches(0.8), clr, num, tsz=22)
    tb(sl, x + Inches(0.1), Inches(2.9), Inches(1.9), Inches(0.8),
       label, sz=16, c=MAIN_DARK, b=True, a=PP_ALIGN.CENTER)
    if num != "5":
        arrow_right(sl, x + Inches(2.15), Inches(3.1), Inches(0.35), clr)
    x += Inches(2.5)

# Bottom: two areas side by side
# Payment methods
rrect(sl, Inches(0.5), Inches(5.0), Inches(6.0), Inches(2.1), WHITE, border_c=BORDER)
tb(sl, Inches(0.7), Inches(5.15), Inches(5.6), Inches(0.35),
   "選べる決済方法", sz=14, c=MAIN_DARK, b=True)
pay_items = [
    ("クレジットカード", MAIN, "おすすめ"),
    ("掛売（請求書払い）", MUTED, "法人向け"),
    ("代金引換", MUTED, ""),
]
py = Inches(5.6)
for name, clr, note in pay_items:
    rect(sl, Inches(0.7), py, Inches(0.08), Inches(0.35), clr)
    tb(sl, Inches(1.0), py + Inches(0.02), Inches(3), Inches(0.3),
       name, sz=13, c=MAIN_DARK, b=(clr == MAIN))
    if note:
        pill(sl, Inches(4.0), py + Inches(0.02), Inches(1.2), Inches(0.3),
             note, bg_c=HIGHLIGHT if note == "おすすめ" else SOFT,
             fg=MAIN_DARK, sz=10, b=False)
    py += Inches(0.42)

# Receive methods
rrect(sl, Inches(6.8), Inches(5.0), Inches(6.0), Inches(2.1), WHITE, border_c=BORDER)
tb(sl, Inches(7.0), Inches(5.15), Inches(5.6), Inches(0.35),
   "選べる受け取り方法", sz=14, c=MAIN_DARK, b=True)
recv_items = [
    ("配送", BLUE_ACC, "ご自宅・店舗へお届け"),
    ("畑で受け取り", MAIN, "送料無料・生産者と直接会える"),
]
ry = Inches(5.65)
for name, clr, desc in recv_items:
    rect(sl, Inches(7.0), ry, Inches(0.08), Inches(0.5), clr)
    tb(sl, Inches(7.3), ry + Inches(0.02), Inches(2.5), Inches(0.25),
       name, sz=13, c=MAIN_DARK, b=True)
    tb(sl, Inches(7.3), ry + Inches(0.28), Inches(4.5), Inches(0.22),
       desc, sz=11, c=MUTED)
    ry += Inches(0.6)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 15–18 — STEP 1–4 (screenshot placeholders)
# ═══════════════════════════════════════════════════════════════
step_slides = [
    ("STEP 1", "商品を探す",
     "トップページや商品一覧から旬の野菜を探せます",
     [("商品一覧ページ", Inches(0.8)), ("商品詳細ページ", Inches(7.0))],
     ["カテゴリ・キーワードで絞り込み", "生産者情報・栽培方法を確認",
      "価格・内容量・配送方法をチェック"]),
    ("STEP 2", "カートに入れる",
     "数量を選んで「カートに追加」。複数の商品をまとめて購入できます",
     [("カート画面", Inches(6.8))],
     ["数量を選んで「カートに追加」", "規格（バリエーション）の選択",
      "カート内で数量変更・削除", "複数出品者の商品もまとめてOK"]),
    ("STEP 3", "購入手続き",
     "受け取り方法・配送先・お届け日時・支払い方法を選択",
     [("チェックアウト画面", Inches(0.8))],
     ["受け取り方法：配送 or 畑受け取り", "配送先の住所を入力 or 選択",
      "お届け日時を選ぶ", "支払い方法を選ぶ", "内容を確認して「注文確定」"]),
    ("STEP 4", "注文確定・確認メール",
     "注文完了画面が表示され、メールで注文確認が届きます",
     [("注文確認画面", Inches(0.8)), ("注文完了画面", Inches(7.0))],
     ["注文番号の発行", "確認メールの自動送信", "ダッシュボードに注文履歴が追加"]),
]

for step_label, title, desc, screenshots, points in step_slides:
    sl = new_slide(); bg(sl, WHITE); stripe(sl)
    pill(sl, Inches(0.8), Inches(0.4), Inches(1.3), Inches(0.38),
         step_label, bg_c=MAIN, sz=12)
    tb(sl, Inches(2.3), Inches(0.33), Inches(6), Inches(0.5),
       title, sz=26, c=MAIN_DARK, b=True)
    tb(sl, Inches(0.8), Inches(0.95), Inches(8), Inches(0.4),
       desc, sz=14, c=MUTED)

    if len(screenshots) == 2:
        ss_placeholder(sl, screenshots[0][1], Inches(1.6), Inches(5.8), Inches(5.4),
                       screenshots[0][0])
        ss_placeholder(sl, screenshots[1][1], Inches(1.6), Inches(5.8), Inches(5.4),
                       screenshots[1][0])
    elif len(screenshots) == 1:
        if step_label in ("STEP 2", "STEP 3"):
            # Points on left/right, screenshot on other side
            if step_label == "STEP 2":
                # Points left, SS right
                y = Inches(1.8)
                for i, p in enumerate(points):
                    clabel(sl, Inches(1.0), y, Inches(0.3), ACCENT,
                           str(i+1), tsz=11, tc=WHITE)
                    tb(sl, Inches(1.5), y + Inches(0.02), Inches(4.8), Inches(0.35),
                       p, sz=14, c=TEXT_CLR)
                    y += Inches(0.5)
                ss_placeholder(sl, screenshots[0][1], Inches(1.5),
                               Inches(5.8), Inches(5.5), screenshots[0][0])
            else:
                # SS left, points right
                ss_placeholder(sl, screenshots[0][1], Inches(1.5),
                               Inches(6.2), Inches(5.5), screenshots[0][0])
                y = Inches(1.8)
                for i, p in enumerate(points):
                    clabel(sl, Inches(7.5), y, Inches(0.35), MAIN,
                           str(i+1), tsz=12, tc=WHITE)
                    tb(sl, Inches(8.1), y + Inches(0.04), Inches(4.5), Inches(0.35),
                       p, sz=14, c=TEXT_CLR)
                    y += Inches(0.55)
    snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 19 — PAYMENT: Comparison Table (credit card emphasized)
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

heading(sl, "選べる決済方法", "ご利用状況に合わせてお選びいただけます")

headers = ["", "クレジットカード", "掛売（請求書払い）", "代金引換"]
rows = [
    ("お支払いタイミング", "注文時に即時決済", "月末締め翌月末払い", "商品到着時"),
    ("請求書", "不要", "月次で自動発行（PDF）", "不要"),
    ("対象", "個人・法人", "法人のみ", "個人・法人"),
    ("セキュリティ", "Stripe社の安全な基盤", "与信限度額で管理", "—"),
]

col_w = Inches(3.2)
col_x = [Inches(1.0), Inches(1.0) + Inches(2.5),
         Inches(1.0) + Inches(2.5) + col_w,
         Inches(1.0) + Inches(2.5) + col_w * 2]

# Header row — credit card gets prominent MAIN color
hy = Inches(2.0)
rect(sl, col_x[1], hy, col_w, Inches(0.55), MAIN)
rect(sl, col_x[2], hy, col_w, Inches(0.55), MUTED)
rect(sl, col_x[3], hy, col_w, Inches(0.55), MUTED)
# "おすすめ" badge on credit card
pill(sl, col_x[1] + Inches(0.9), hy - Inches(0.3), Inches(1.4), Inches(0.3),
     "おすすめ", bg_c=GOLD, fg=WHITE, sz=10)
for i in range(1, 4):
    tb(sl, col_x[i], hy + Inches(0.1), col_w, Inches(0.4),
       headers[i], sz=13, c=WHITE, b=True, a=PP_ALIGN.CENTER)

# Data rows
ry = Inches(2.6)
for i, (label, *vals) in enumerate(rows):
    row_bg = SOFT if i % 2 == 0 else WHITE
    rect(sl, col_x[0], ry, Inches(2.5), Inches(0.65), row_bg)
    for j in range(3):
        rect(sl, col_x[j+1], ry, col_w, Inches(0.65), row_bg)

    tb(sl, col_x[0] + Inches(0.15), ry + Inches(0.15), Inches(2.2), Inches(0.4),
       label, sz=12, c=MAIN_DARK, b=True)
    for j, val in enumerate(vals):
        tb(sl, col_x[j+1] + Inches(0.15), ry + Inches(0.15), col_w - Inches(0.3), Inches(0.4),
           val, sz=12, c=TEXT_CLR, a=PP_ALIGN.CENTER)
    ry += Inches(0.65)

# Highlight: credit card recommendation
rrect(sl, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.2), HIGHLIGHT)
tb(sl, Inches(1.3), Inches(5.6), Inches(10.7), Inches(0.4),
   "まずは「クレジットカード決済」がおすすめです",
   sz=16, c=MAIN_DARK, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(1.3), Inches(6.05), Inches(10.7), Inches(0.5),
   "Stripe社の安全な決済基盤を利用。カード情報はセッツマルシェには保存されません。\n法人のお客様は、ご希望に応じて掛売（請求書払い）もご利用いただけます。",
   sz=13, c=MUTED, a=PP_ALIGN.CENTER, ls=22)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 20 — CREDIT CARD PAYMENT FLOW + screenshot
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

pill(sl, Inches(0.8), Inches(0.4), Inches(1.8), Inches(0.38),
     "おすすめ", bg_c=MAIN, sz=12)
tb(sl, Inches(2.8), Inches(0.33), Inches(7), Inches(0.5),
   "クレジットカード決済の手順", sz=26, c=MAIN_DARK, b=True)

cc_steps = [
    ("1", "購入手続き画面で「クレジットカード」を選択"),
    ("2", "カード番号・有効期限・セキュリティコードを入力"),
    ("3", "「注文を確定する」ボタンで即時決済"),
    ("4", "注文完了。確認メールが届きます"),
]
y = Inches(1.3)
for num, text in cc_steps:
    clabel(sl, Inches(1.0), y, Inches(0.4), MAIN, num, tsz=15)
    tb(sl, Inches(1.7), y + Inches(0.06), Inches(4.8), Inches(0.35),
       text, sz=15, c=TEXT_CLR)
    if num != "4":
        rect(sl, Inches(1.19), y + Inches(0.42), Inches(0.02), Inches(0.35), ACCENT)
    y += Inches(0.68)

# Safety note
rrect(sl, Inches(0.8), Inches(4.5), Inches(5.5), Inches(1.2), SOFT)
tb(sl, Inches(1.0), Inches(4.6), Inches(5.1), Inches(0.3),
   "安心のセキュリティ", sz=14, c=MAIN_DARK, b=True)
tb(sl, Inches(1.0), Inches(4.95), Inches(5.1), Inches(0.6),
   "決済はStripe社の基盤を利用しており、\nカード情報はセッツマルシェに保存されません。",
   sz=12, c=MUTED, ls=20)

# Note about other options
tb(sl, Inches(0.8), Inches(6.0), Inches(5.5), Inches(0.5),
   "※ 法人のお客様は掛売（請求書払い）もご利用可能です",
   sz=11, c=MUTED)

ss_placeholder(sl, Inches(6.8), Inches(1.1), Inches(5.8), Inches(6.0),
               "チェックアウト画面（決済部分）")

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 21 — DELIVERY STATUS + screenshot
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

pill(sl, Inches(0.8), Inches(0.4), Inches(1.5), Inches(0.38),
     "注文後", bg_c=MAIN, sz=12)
tb(sl, Inches(2.5), Inches(0.33), Inches(7), Inches(0.5),
   "配送状況の確認", sz=26, c=MAIN_DARK, b=True)
tb(sl, Inches(0.8), Inches(0.95), Inches(8), Inches(0.4),
   "ダッシュボードからリアルタイムで確認できます", sz=14, c=MUTED)

ss_placeholder(sl, Inches(0.8), Inches(1.5), Inches(6.2), Inches(5.5),
               "注文詳細 / 配送ステータス画面")

statuses = [
    ("出荷準備中", "商品を準備しています", ACCENT),
    ("お届け中", "配送業者に渡しました", BLUE_ACC),
    ("配達完了", "お届けが完了しました", MAIN),
]
y = Inches(1.8)
for label, desc, color in statuses:
    c = rrect(sl, Inches(7.4), y, Inches(5.0), Inches(0.7), WHITE, border_c=BORDER)
    rect(sl, Inches(7.4), y, Inches(0.08), Inches(0.7), color)
    tb(sl, Inches(7.7), y + Inches(0.1), Inches(2.5), Inches(0.3),
       label, sz=14, c=MAIN_DARK, b=True)
    tb(sl, Inches(7.7), y + Inches(0.38), Inches(3.5), Inches(0.25),
       desc, sz=12, c=MUTED)
    if label != "配達完了":
        # connector line
        rect(sl, Inches(7.43), y + Inches(0.7), Inches(0.02), Inches(0.2), BORDER)
    y += Inches(0.9)

pill(sl, Inches(7.4), Inches(4.8), Inches(5.0), Inches(0.65),
     "注文が入るとメール通知が届き\nダッシュボードからもいつでも確認できます",
     bg_c=HIGHLIGHT, fg=MAIN_DARK, sz=12, b=True)

# Pickup status on same slide
tb(sl, Inches(7.4), Inches(5.8), Inches(5.0), Inches(0.35),
   "畑受け取りの場合", sz=13, c=MAIN_DARK, b=True)

p_statuses = [("受け取り準備中", ACCENT), ("準備完了", BLUE_ACC), ("受け渡し完了", MAIN)]
px = Inches(7.4)
for lbl, clr in p_statuses:
    pill(sl, px, Inches(6.25), Inches(1.55), Inches(0.35), lbl,
         bg_c=clr, fg=WHITE, sz=10)
    if lbl != "受け渡し完了":
        tb(sl, px + Inches(1.55), Inches(6.25), Inches(0.15), Inches(0.35),
           "→", sz=12, c=MUTED, a=PP_ALIGN.CENTER)
    px += Inches(1.7)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 22 — FARM PICKUP (comparison diagram)
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

pill(sl, Inches(0.8), Inches(0.4), Inches(2.8), Inches(0.38),
     "セッツマルシェならでは", bg_c=MAIN, sz=12)
tb(sl, Inches(3.8), Inches(0.33), Inches(7), Inches(0.5),
   "畑での受け取り", sz=26, c=MAIN_DARK, b=True)

# Clean comparison table with clear columns
# Full-width table background
rect(sl, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8), WHITE)

# Column headers
lbl_w = Inches(2.0)
left_w = Inches(4.8)
right_w = Inches(5.0)
lbl_x = Inches(0.5)
left_x = Inches(2.5)
right_x = Inches(7.6)

rect(sl, lbl_x, Inches(1.2), lbl_w, Inches(0.6), SOFT)
rect(sl, left_x, Inches(1.2), left_w, Inches(0.6), RGBColor(0xE0, 0xEA, 0xF0))
rect(sl, right_x, Inches(1.2), right_w, Inches(0.6), HIGHLIGHT)

tb(sl, left_x, Inches(1.25), left_w, Inches(0.5),
   "配送", sz=16, c=BLUE_ACC, b=True, a=PP_ALIGN.CENTER)
tb(sl, right_x, Inches(1.25), right_w, Inches(0.5),
   "畑で受け取り", sz=16, c=MAIN, b=True, a=PP_ALIGN.CENTER)

comp_rows = [
    ("送料", "地域・重量に応じて発生", "無料"),
    ("鮮度", "梱包・配送を経てお届け", "収穫したてを直接受け取り"),
    ("手間", "住所入力・日時指定が必要", "「畑受け取り」を選ぶだけ"),
    ("体験", "通常のネット通販と同様", "生産者と直接会える体験"),
]
ry = Inches(1.85)
for i, (label, left_val, right_val) in enumerate(comp_rows):
    row_h = Inches(0.85)
    # Alternating row backgrounds
    even_bg = RGBColor(0xFA, 0xFA, 0xF8) if i % 2 == 0 else WHITE
    rect(sl, lbl_x, ry, lbl_w, row_h, SOFT)
    rect(sl, left_x, ry, left_w, row_h, even_bg)
    rect(sl, right_x, ry, right_w, row_h, even_bg)

    # Row borders
    rect(sl, lbl_x, ry + row_h - Inches(0.01), Inches(12.1), Inches(0.01), BORDER)

    tb(sl, lbl_x + Inches(0.2), ry + Inches(0.22), lbl_w - Inches(0.4), Inches(0.4),
       label, sz=13, c=MAIN_DARK, b=True, a=PP_ALIGN.CENTER)
    tb(sl, left_x + Inches(0.2), ry + Inches(0.22), left_w - Inches(0.4), Inches(0.4),
       left_val, sz=13, c=MUTED, a=PP_ALIGN.CENTER)
    tb(sl, right_x + Inches(0.2), ry + Inches(0.22), right_w - Inches(0.4), Inches(0.4),
       right_val, sz=13, c=MAIN_DARK, b=True, a=PP_ALIGN.CENTER)
    ry += row_h

# Highlight for farm pickup
rrect(sl, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.7), HIGHLIGHT)
tb(sl, Inches(0.7), Inches(5.8), Inches(11.9), Inches(0.5),
   "畑で受け取り ＝ 生産者の顔が見える安心感 × 送料ゼロのお得さ。セッツマルシェならではの体験です。",
   sz=14, c=MAIN_DARK, b=True, a=PP_ALIGN.CENTER)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 23 — REPEAT ORDER
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

heading(sl, "便利な機能 ①  リピート注文", "過去の注文をワンクリックで再注文")

# Flow diagram
flow_items = [
    ("注文履歴を\n開く", MAIN),
    ("「この注文を\nもう一度」", GOLD),
    ("カートに\n自動追加", ACCENT),
    ("内容確認\n→ 注文確定", MAIN),
]
x = Inches(0.8)
for label, clr in flow_items:
    rrect(sl, x, Inches(1.8), Inches(2.5), Inches(1.6), clr)
    tb(sl, x + Inches(0.1), Inches(2.2), Inches(2.3), Inches(0.8),
       label, sz=15, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    if label != flow_items[-1][0]:
        arrow_right(sl, x + Inches(2.55), Inches(2.4), Inches(0.4), ACCENT)
    x += Inches(2.95)

# Screenshot below
ss_placeholder(sl, Inches(0.8), Inches(3.8), Inches(11.7), Inches(3.2),
               "注文詳細画面 —「この注文をもう一度」ボタン")

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 24 — ORDER TEMPLATE
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, WHITE); stripe(sl)

heading(sl, "便利な機能 ②  注文テンプレート", "「いつもの注文」をワンクリックで")

# Two ways to create
tb(sl, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.35),
   "テンプレートの作り方（2通り）", sz=15, c=MAIN_DARK, b=True)

ways = [
    ("カートから保存", "カート画面で「テンプレートとして保存」"),
    ("注文履歴から保存", "注文詳細で「テンプレートとして保存」"),
]
wy = Inches(2.1)
for title, desc in ways:
    rrect(sl, Inches(0.8), wy, Inches(5.5), Inches(0.7), SOFT)
    tb(sl, Inches(1.0), wy + Inches(0.1), Inches(2), Inches(0.3),
       title, sz=13, c=MAIN_DARK, b=True)
    tb(sl, Inches(1.0), wy + Inches(0.38), Inches(5), Inches(0.25),
       desc, sz=12, c=MUTED)
    wy += Inches(0.85)

# Usage flow
tb(sl, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.35),
   "テンプレートの使い方", sz=15, c=MAIN_DARK, b=True)

use_steps = [
    "「注文テンプレート」を開く",
    "使いたいテンプレートの「カートに追加」",
    "全商品が自動でカートに追加",
    "内容確認 → 購入手続きへ",
]
uy = Inches(4.5)
for i, s in enumerate(use_steps):
    clabel(sl, Inches(1.0), uy, Inches(0.3), MAIN, str(i+1), tsz=11, tc=WHITE)
    tb(sl, Inches(1.5), uy + Inches(0.02), Inches(4.8), Inches(0.3),
       s, sz=13, c=TEXT_CLR)
    uy += Inches(0.42)

# Screenshot right
ss_placeholder(sl, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5),
               "注文テンプレート一覧")

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 25 — B2B FEATURES
# ═══════════════════════════════════════════════════════════════
sl = new_slide(); bg(sl, SOFT); stripe(sl)

heading(sl, "法人向けの安心機能", "組織での利用に対応した仕組みをご用意しています")

b2b = [
    (MAIN, "¥", "顧客別価格",
     "お取引状況に応じた特別価格を設定。\nログインすると自動で適用されます。",
     ["取引量に応じたボリュームディスカウント", "バリエーションごとに個別設定可能"]),
    (ACCENT, "✓", "購買承認フロー",
     "担当者が注文 → 店長が承認。\n組織に合わせた承認ステップを設定できます。",
     ["多段階の承認にも対応", "承認不要の設定も可能"]),
    (GOLD, "⇔", "組織メンバー管理",
     "注文担当・承認者・経理など、\nメンバーごとに権限を設定できます。",
     ["注文担当 / 承認者 / 経理 / 管理者", "最低1名の管理者が必須"]),
]
x = Inches(0.5)
for clr, icon, title, desc, points in b2b:
    card = rrect(sl, x, Inches(1.7), Inches(3.9), Inches(5.2), WHITE, border_c=BORDER)
    rect(sl, x, Inches(1.7), Inches(3.9), Inches(0.07), clr)
    clabel(sl, x + Inches(1.35), Inches(2.1), Inches(1.2), clr, icon, tsz=30, tc=WHITE)
    tb(sl, x + Inches(0.3), Inches(3.5), Inches(3.3), Inches(0.4),
       title, sz=17, c=MAIN_DARK, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.3), Inches(4.0), Inches(3.3), Inches(0.9),
       desc, sz=12, c=MUTED, a=PP_ALIGN.CENTER, ls=20)
    py = Inches(5.1)
    for pt in points:
        tb(sl, x + Inches(0.5), py, Inches(3), Inches(0.3),
           "・" + pt, sz=11, c=MAIN_DARK)
        py += Inches(0.3)
    x += Inches(4.2)

snum(sl)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 26 — CLOSING
# ═══════════════════════════════════════════════════════════════
sl = new_slide()
img_fill(sl, img_path("slide_closing.png"))
rect(sl, Emu(0), Emu(0), SW, SH, DARK_BG, alpha=60)

tb(sl, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
   "地域でつくり、地域で食べる。\nその循環を、ごいっしょに。",
   sz=36, c=WHITE, b=True, a=PP_ALIGN.CENTER, f=SERIF, ls=54)

tb(sl, Inches(2.5), Inches(3.4), Inches(8), Inches(0.8),
   "セッツマルシェは、飲食店のみなさまと一緒に育てていくサービスです。\nご意見・ご要望をぜひお聞かせください。",
   sz=15, c=WHITE, a=PP_ALIGN.CENTER, ls=26)

info = rrect(sl, Inches(3.5), Inches(4.7), Inches(6.3), Inches(1.5),
             RGBColor(0x3A, 0x55, 0x45))
mtb(sl, Inches(3.7), Inches(4.85), Inches(5.9), Inches(1.2),
    [
        ("セッツマルシェ（渡辺ファーム）", 16, WHITE, True),
        ("", 8, WHITE, False),
        ("お問い合わせ：サイト内「お問い合わせフォーム」より", 13, HIGHLIGHT, False),
    ],
    a=PP_ALIGN.CENTER,
)


# ═══════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"✅ Saved: {OUT}")
print(f"   Slides: {slide_n[0]}")
