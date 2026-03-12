#!/usr/bin/env python3
"""セッツマルシェ プレゼンテーション .pptx ビルダー"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── paths ──────────────────────────────────────────
ASSETS = "/Users/koya1104/.cursor/projects/Users-koya1104-Desktop/assets"
OUT    = "/Users/koya1104/Desktop/新・今日の食卓/z_document/presentation/セッツマルシェ_ご紹介資料.pptx"

# ── brand tokens ───────────────────────────────────
MAIN       = RGBColor(0x4C, 0x6B, 0x5C)
MAIN_DARK  = RGBColor(0x35, 0x5D, 0x4A)
ACCENT     = RGBColor(0xA3, 0xC9, 0xA8)
HIGHLIGHT  = RGBColor(0xD6, 0xEA, 0xDF)
SOFT       = RGBColor(0xF2, 0xF7, 0xF4)
BG         = RGBColor(0xF9, 0xF9, 0xF6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_COLOR  = RGBColor(0x2B, 0x2B, 0x2B)
MUTED      = RGBColor(0x6B, 0x72, 0x80)
WARN_RED   = RGBColor(0xE8, 0x92, 0x7C)

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)
FONT_SANS = "Noto Sans JP"
FONT_SERIF = "Noto Serif JP"

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank layout

# ── helpers ────────────────────────────────────────

def img(name):
    return os.path.join(ASSETS, name)

def add_bg_rect(slide, color, left=0, top=0, w=None, h=None, alpha=None):
    """alpha: 0–100 (0=fully transparent, 100=fully opaque). None = fully opaque."""
    from pptx.oxml.ns import qn
    from lxml import etree
    w = w or SLIDE_W
    h = h or SLIDE_H
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    if alpha is not None:
        sp_pr = shp._element.find(qn('a:solidFill'), shp._element.nsmap)
        if sp_pr is None:
            sp_pr_parent = shp._element.find(qn('p:spPr'), shp._element.nsmap)
            if sp_pr_parent is None:
                sp_pr_parent = shp._element.spPr
            solid_fill = sp_pr_parent.find(qn('a:solidFill'))
            if solid_fill is not None:
                for child in solid_fill:
                    a_el = child.find(qn('a:alpha'))
                    if a_el is None:
                        a_el = etree.SubElement(child, qn('a:alpha'))
                    a_el.set('val', str(alpha * 1000))
    return shp

def add_gradient_bg(slide, c1, c2):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill = shp.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0
    shp.line.fill.background()
    shp.rotation = 0
    return shp

def add_bottom_stripe(slide):
    h = Inches(0.08)
    colors = [MAIN, ACCENT, HIGHLIGHT]
    w_each = SLIDE_W // 3
    for i, c in enumerate(colors):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(w_each) * i), SLIDE_H - h, Emu(int(w_each) + 1), h
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = c
        shp.line.fill.background()

def add_slide_num(slide, num, total, color=MUTED):
    txBox = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.5), Inches(1), Inches(0.35)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.RIGHT

def add_text(slide, left, top, w, h, text, font_size=14,
             bold=False, color=TEXT_COLOR, align=PP_ALIGN.LEFT,
             font_name=FONT_SANS, line_spacing=1.3, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.font.name = font_name
    p.line_spacing = Pt(int(font_size * line_spacing))
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox

def add_multiline(slide, left, top, w, h, lines, font_size=14,
                  color=TEXT_COLOR, align=PP_ALIGN.LEFT, font_name=FONT_SANS,
                  line_spacing=1.5, bold=False):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        p.font.bold = bold
        p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox

def add_rounded_rect(slide, left, top, w, h, fill_color, border_color=None, radius=Inches(0.15)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    return shp

def add_circle(slide, left, top, size, fill_color):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

def add_screenshot_placeholder(slide, left, top, w, h, label):
    shp = add_rounded_rect(slide, left, top, w, h, SOFT, ACCENT)
    shp.line.dash_style = 4  # dash
    add_text(slide, left, top + h // 2 - Inches(0.3), w, Inches(0.3),
             "[ Screenshot ]", font_size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + h // 2 + Inches(0.05), w, Inches(0.3),
             label, font_size=10, color=MAIN_DARK, align=PP_ALIGN.CENTER, bold=True)

TOTAL = 20

# ================================================================
# SLIDE 1 — Title
# ================================================================
s = prs.slides.add_slide(BLANK)
add_gradient_bg(s, MAIN, MAIN_DARK)

s.shapes.add_picture(img("slide_hero.png"),
                     Inches(0), Inches(0), SLIDE_W, SLIDE_H)
add_bg_rect(s, MAIN_DARK, Inches(0), Inches(0), SLIDE_W, SLIDE_H, alpha=60)

add_text(s, Inches(0), Inches(1.5), SLIDE_W, Inches(0.4),
         "2026.03.24 ご紹介資料", font_size=13, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(2.3), SLIDE_W, Inches(1.2),
         "セッツマルシェ", font_size=52, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF)
add_text(s, Inches(0), Inches(3.6), SLIDE_W, Inches(0.7),
         "まちと畑を、ひとつの食卓につなぐ。", font_size=22,
         color=WHITE, align=PP_ALIGN.CENTER, font_name=FONT_SERIF)
add_text(s, Inches(0), Inches(5.8), SLIDE_W, Inches(0.4),
         "渡辺ファーム ／ 摂津市商工会", font_size=12,
         color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 2 — Agenda
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_text(s, Inches(0), Inches(0.8), SLIDE_W, Inches(0.6),
         "本日の内容", font_size=32, bold=True, color=MAIN_DARK,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF)

agenda_items = [
    ("1", "セッツマルシェとは", "背景・コンセプト・目指す姿"),
    ("2", "わたしたちの柱", "小さな声を拾う ／ まちの食を支える ／ 本当の美味しさを届ける"),
    ("3", "野菜の出どころ", "渡辺ファーム・小規模農家・市民農園"),
    ("4", "購入フローのご紹介", "購入手順・決済・配送・畑受け取り"),
    ("5", "便利な機能と今後", "リピート注文・テンプレート・掛売"),
]
y_start = Inches(2.0)
for i, (num, title, sub) in enumerate(agenda_items):
    y = y_start + Inches(i * 0.95)
    add_circle(s, Inches(4.0), y, Inches(0.5), MAIN)
    add_text(s, Inches(4.0), y + Inches(0.08), Inches(0.5), Inches(0.4),
             num, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.8), y + Inches(0.02), Inches(5), Inches(0.35),
             title, font_size=17, bold=True, color=MAIN_DARK)
    add_text(s, Inches(4.8), y + Inches(0.38), Inches(5), Inches(0.3),
             sub, font_size=11, color=MUTED)

add_slide_num(s, 2, TOTAL)

# ================================================================
# SLIDE 3 — Background / Problem
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_text(s, Inches(0), Inches(0.7), SLIDE_W, Inches(0.8),
         "地域の「もったいない」を\n「おいしい」に変えたい", font_size=30,
         bold=True, color=MAIN_DARK, align=PP_ALIGN.CENTER,
         font_name=FONT_SERIF)

card_w = Inches(3.5)
card_h = Inches(3.8)
gap = Inches(0.4)
total_w = card_w * 3 + gap * 2
start_x = (SLIDE_W - total_w) // 2
y_card = Inches(2.2)

problems = [
    ("小規模農家の課題", "採れすぎた野菜や\n小さな畑の収穫が\n販路に乗らず\n廃棄されてしまう", img("pillar_small_voices.png")),
    ("飲食店の課題",     "新鮮な地場野菜を\n安定的に仕入れる\nルートが\n限られている",   img("pillar_support_dining.png")),
    ("地域の課題",       "作る人・届ける人\n味わう人の\nつながりが\n見えにくくなっている", img("pillar_true_delicious.png")),
]

for i, (title, desc, image) in enumerate(problems):
    x = Emu(int(start_x) + int(Emu(int(card_w) + int(gap))) * i)
    add_rounded_rect(s, x, y_card, card_w, card_h, BG, RGBColor(0xE5, 0xE7, 0xEB))
    accent_bar = slide_shapes = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y_card, card_w, Inches(0.06))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = WARN_RED
    accent_bar.line.fill.background()
    pic_w = Inches(2.4)
    pic_h = Inches(1.4)
    pic_x = x + (card_w - pic_w) // 2
    s.shapes.add_picture(image, pic_x, y_card + Inches(0.25), pic_w, pic_h)
    add_text(s, x + Inches(0.3), y_card + Inches(1.85), card_w - Inches(0.6), Inches(0.4),
             title, font_size=16, bold=True, color=MAIN_DARK, align=PP_ALIGN.CENTER)
    add_multiline(s, x + Inches(0.3), y_card + Inches(2.35), card_w - Inches(0.6), Inches(1.4),
                  desc.split('\n'), font_size=13, color=MUTED, align=PP_ALIGN.CENTER,
                  line_spacing=1.6)

highlight_rect = add_rounded_rect(s, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.65), HIGHLIGHT)
add_text(s, Inches(2.5), Inches(6.35), Inches(8.3), Inches(0.55),
         "セッツマルシェは、これらの課題をやさしく解決する地産地消の仕組みです。",
         font_size=14, color=MAIN_DARK, align=PP_ALIGN.CENTER, bold=True)
add_slide_num(s, 3, TOTAL)

# ================================================================
# SLIDE 4 — What is セッツマルシェ
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)

s.shapes.add_picture(img("sources_diagram.png"),
                     Inches(0), Inches(0), SLIDE_W, SLIDE_H)
add_bg_rect(s, WHITE, Inches(0), Inches(0), SLIDE_W, SLIDE_H, alpha=30)

add_text(s, Inches(0), Inches(0.6), SLIDE_W, Inches(0.5),
         "セッツマルシェとは", font_size=32, bold=True, color=MAIN_DARK,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF)
add_text(s, Inches(2), Inches(1.3), Inches(9.3), Inches(0.6),
         "市民農園・小規模農家の野菜を、地域の飲食店へ届けるオンラインマルシェ",
         font_size=16, color=MUTED, align=PP_ALIGN.CENTER)

s.shapes.add_picture(img("sources_diagram.png"),
                     Inches(1.5), Inches(2.1), Inches(10.3), Inches(5.0))

add_bottom_stripe(s)
add_slide_num(s, 4, TOTAL)

# ================================================================
# SLIDE 5 — 3 Pillars (modified: removed やさしい流通, added 本当の美味しさ)
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, SOFT)
add_bottom_stripe(s)

add_text(s, Inches(0), Inches(0.6), SLIDE_W, Inches(0.5),
         "わたしたちの 3つの柱", font_size=32, bold=True, color=MAIN_DARK,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF)

pillars = [
    (img("pillar_small_voices.png"),  "小さな声を拾う",
     "家庭菜園や小規模農家の収穫を活かし、\n持続可能な販売機会に変えます。"),
    (img("pillar_support_dining.png"), "まちの食を支える",
     "地域の飲食店や家庭に、\n新鮮で確かな食材を安定して届けます。"),
    (img("pillar_true_delicious.png"), "本当の美味しさを届ける",
     "畑から食卓へ、最短ルートで届ける鮮度。\nこの距離の近さが、本物の美味しさを生みます。"),
]

p_card_w = Inches(3.6)
p_card_h = Inches(5.0)
p_gap = Inches(0.35)
p_total = p_card_w * 3 + p_gap * 2
p_start_x = (SLIDE_W - p_total) // 2
p_y = Inches(1.6)

for i, (image, title, desc) in enumerate(pillars):
    x = Emu(int(p_start_x) + int(Emu(int(p_card_w) + int(p_gap))) * i)
    add_rounded_rect(s, x, p_y, p_card_w, p_card_h, WHITE, RGBColor(0xE5, 0xE7, 0xEB))
    top_bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, p_y, p_card_w, Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = [MAIN, ACCENT, MAIN_DARK][i]
    top_bar.line.fill.background()
    pic_w = Inches(3.0)
    pic_h = Inches(1.8)
    s.shapes.add_picture(image, x + (p_card_w - pic_w) // 2, p_y + Inches(0.3), pic_w, pic_h)
    add_text(s, x + Inches(0.3), p_y + Inches(2.3), p_card_w - Inches(0.6), Inches(0.4),
             title, font_size=18, bold=True, color=MAIN_DARK, align=PP_ALIGN.CENTER)
    add_multiline(s, x + Inches(0.3), p_y + Inches(2.9), p_card_w - Inches(0.6), Inches(1.5),
                  desc.split('\n'), font_size=12, color=MUTED, align=PP_ALIGN.CENTER,
                  line_spacing=1.7)

add_slide_num(s, 5, TOTAL)

# ================================================================
# SLIDE 6 — Sources of Vegetables
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_text(s, Inches(0), Inches(0.6), SLIDE_W, Inches(0.5),
         "取り扱う野菜の出どころ", font_size=32, bold=True, color=MAIN_DARK,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF)

s.shapes.add_picture(img("sources_diagram.png"),
                     Inches(0.8), Inches(1.5), Inches(6.5), Inches(4.5))

src_items = [
    ("渡辺ファーム", "運営元の農場。安心・安全な栽培で、\n主力となる野菜を安定供給します。", MAIN),
    ("小規模農家さん", "地域で丁寧に野菜を育てる農家さん。\n販路を共有し、収穫を無駄にしません。", ACCENT),
    ("市民農園の方々", "家庭菜園の余剰も、飲食店の一皿に。\n小さな生産にも価値を見出します。", HIGHLIGHT),
]

right_x = Inches(7.8)
for i, (title, desc, bar_color) in enumerate(src_items):
    y = Inches(1.8) + Inches(i * 1.7)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, y, Inches(0.08), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    add_text(s, right_x + Inches(0.3), y, Inches(4.5), Inches(0.35),
             title, font_size=17, bold=True, color=MAIN_DARK)
    add_multiline(s, right_x + Inches(0.3), y + Inches(0.45), Inches(4.5), Inches(0.8),
                  desc.split('\n'), font_size=12, color=MUTED, line_spacing=1.6)

add_slide_num(s, 6, TOTAL)

# ================================================================
# SLIDE 7 — Vision
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, SOFT)

s.shapes.add_picture(img("circulation_vision.png"),
                     Inches(0), Inches(0), SLIDE_W, SLIDE_H)
add_bg_rect(s, SOFT, alpha=45)

add_text(s, Inches(0), Inches(2.0), SLIDE_W, Inches(1.8),
         "地域でつくり、\n地域で味わい、\n地域でつながる。",
         font_size=36, bold=True, color=MAIN_DARK,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF, line_spacing=1.6)
add_text(s, Inches(2.5), Inches(4.8), Inches(8.3), Inches(0.9),
         "この循環の輪が、まちの豊かさを育んでいく。\nそれが、私たち「セッツマルシェ」の願いです。",
         font_size=16, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.7)

add_bottom_stripe(s)
add_slide_num(s, 7, TOTAL)

# ================================================================
# SLIDE 8 — Part 2 Section Divider
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, SOFT)
add_bottom_stripe(s)

label_w = Inches(1.5)
label_h = Inches(0.4)
label_x = (SLIDE_W - label_w) // 2
add_rounded_rect(s, label_x, Inches(2.5), label_w, label_h, MAIN)
add_text(s, label_x, Inches(2.52), label_w, label_h,
         "PART 2", font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(0), Inches(3.2), SLIDE_W, Inches(0.7),
         "ご利用方法のご紹介", font_size=34, bold=True, color=MAIN_DARK,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF)
add_text(s, Inches(0), Inches(4.1), SLIDE_W, Inches(0.4),
         "飲食店さま向け — 購入フローと便利な機能", font_size=15,
         color=MUTED, align=PP_ALIGN.CENTER)
add_slide_num(s, 8, TOTAL)

# ================================================================
# SLIDE 9 — Purchase Flow Overview
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_text(s, Inches(0), Inches(0.6), SLIDE_W, Inches(0.5),
         "購入の流れ — 全体像", font_size=30, bold=True, color=MAIN_DARK,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF)

flow_steps = [
    ("商品を探す", MAIN),
    ("カートに入れる", MAIN),
    ("購入手続き", MAIN),
    ("決済", MAIN),
    ("注文確定", MAIN),
]

step_w = Inches(1.8)
step_h = Inches(1.6)
arrow_w = Inches(0.5)
total_flow_w = step_w * 5 + arrow_w * 4
flow_start_x = (SLIDE_W - total_flow_w) // 2
flow_y = Inches(2.2)

for i, (label, color) in enumerate(flow_steps):
    x = Emu(int(flow_start_x) + int(Emu(int(step_w) + int(arrow_w))) * i)
    add_rounded_rect(s, x, flow_y, step_w, step_h, WHITE, color)
    add_circle(s, x + (step_w - Inches(0.6)) // 2, flow_y + Inches(0.2), Inches(0.6), color)
    add_text(s, x + (step_w - Inches(0.6)) // 2, flow_y + Inches(0.28), Inches(0.6), Inches(0.45),
             str(i + 1), font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, flow_y + Inches(0.95), step_w, Inches(0.5),
             label, font_size=14, bold=True, color=MAIN_DARK, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow_x = Emu(int(x) + int(step_w) + int(arrow_w) // 4)
        add_text(s, arrow_x, flow_y + Inches(0.5), Inches(0.3), Inches(0.4),
                 "→", font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

tags = ["クレジットカード", "代金引換", "掛売（請求書払い）", "配送", "畑で受け取り"]
tag_y = Inches(4.6)
tag_total_w = len(tags) * Inches(2.2) + (len(tags) - 1) * Inches(0.2)
tag_start = (SLIDE_W - tag_total_w) // 2
for i, t in enumerate(tags):
    tx = Emu(int(tag_start) + int(Inches(2.4)) * i)
    add_rounded_rect(s, tx, tag_y, Inches(2.1), Inches(0.4), HIGHLIGHT)
    add_text(s, tx, tag_y + Inches(0.05), Inches(2.1), Inches(0.3),
             t, font_size=11, bold=True, color=MAIN_DARK, align=PP_ALIGN.CENTER)

add_slide_num(s, 9, TOTAL)

# ================================================================
# SLIDE 10 — STEP 1: Find Products
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

step_label = add_rounded_rect(s, Inches(0.8), Inches(0.6), Inches(1.2), Inches(0.35), MAIN)
add_text(s, Inches(0.8), Inches(0.62), Inches(1.2), Inches(0.35),
         "STEP 1", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
         "商品を探す", font_size=28, bold=True, color=MAIN_DARK, font_name=FONT_SERIF)
add_text(s, Inches(0.8), Inches(1.7), Inches(5), Inches(0.4),
         "トップページや商品一覧から、旬の野菜を探せます", font_size=13, color=MUTED)

add_screenshot_placeholder(s, Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.3),
                           "商品一覧ページ")
add_screenshot_placeholder(s, Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.3),
                           "商品詳細ページ")
add_slide_num(s, 10, TOTAL)

# ================================================================
# SLIDE 11 — STEP 2: Add to Cart
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_rounded_rect(s, Inches(0.8), Inches(0.6), Inches(1.2), Inches(0.35), MAIN)
add_text(s, Inches(0.8), Inches(0.62), Inches(1.2), Inches(0.35),
         "STEP 2", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
         "カートに入れる", font_size=28, bold=True, color=MAIN_DARK, font_name=FONT_SERIF)
add_text(s, Inches(0.8), Inches(1.7), Inches(5), Inches(0.4),
         "数量を選んで「カートに追加」。複数の商品をまとめて購入できます", font_size=13, color=MUTED)

cart_points = [
    "数量を選んで「カートに追加」",
    "カートアイコンで中身を確認",
    "数量の変更・削除もカート内で可能",
    "規格（バリエーション）がある場合は選択",
]
for i, pt in enumerate(cart_points):
    y = Inches(2.5) + Inches(i * 0.55)
    add_circle(s, Inches(1.0), y, Inches(0.3), ACCENT)
    add_text(s, Inches(1.5), y - Inches(0.02), Inches(4.5), Inches(0.35),
             pt, font_size=13, color=TEXT_COLOR)

add_rounded_rect(s, Inches(0.8), Inches(5.0), Inches(5.5), Inches(0.55), HIGHLIGHT)
add_text(s, Inches(1.0), Inches(5.08), Inches(5.3), Inches(0.45),
         "複数の出品者の商品もまとめてカートに入れられます",
         font_size=12, color=MAIN_DARK, bold=True, align=PP_ALIGN.CENTER)

add_screenshot_placeholder(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.8),
                           "カート画面")
add_slide_num(s, 11, TOTAL)

# ================================================================
# SLIDE 12 — STEP 3: Checkout
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_rounded_rect(s, Inches(0.8), Inches(0.6), Inches(1.2), Inches(0.35), MAIN)
add_text(s, Inches(0.8), Inches(0.62), Inches(1.2), Inches(0.35),
         "STEP 3", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
         "購入手続き", font_size=28, bold=True, color=MAIN_DARK, font_name=FONT_SERIF)
add_text(s, Inches(0.8), Inches(1.7), Inches(9), Inches(0.4),
         "受け取り方法・配送先・お届け日時・支払い方法を選択", font_size=13, color=MUTED)

add_screenshot_placeholder(s, Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.3),
                           "チェックアウト画面")

checkout_steps = [
    ("❶", "受け取り方法を選ぶ（配送 or 畑受け取り）"),
    ("❷", "配送先の住所を入力 or 選択"),
    ("❸", "お届け日時を選ぶ"),
    ("❹", "支払い方法を選ぶ"),
    ("❺", "内容を確認して「注文確定」"),
]
for i, (num, text) in enumerate(checkout_steps):
    y = Inches(2.6) + Inches(i * 0.7)
    add_circle(s, Inches(7.2), y, Inches(0.4), MAIN)
    add_text(s, Inches(7.2), y + Inches(0.06), Inches(0.4), Inches(0.3),
             str(i + 1), font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.8), y + Inches(0.05), Inches(4.8), Inches(0.35),
             text, font_size=13, color=TEXT_COLOR)

add_slide_num(s, 12, TOTAL)

# ================================================================
# SLIDE 13 — STEP 4: Order Confirmation
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_rounded_rect(s, Inches(0.8), Inches(0.6), Inches(1.2), Inches(0.35), MAIN)
add_text(s, Inches(0.8), Inches(0.62), Inches(1.2), Inches(0.35),
         "STEP 4", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
         "注文確定・確認メール", font_size=28, bold=True, color=MAIN_DARK, font_name=FONT_SERIF)
add_text(s, Inches(0.8), Inches(1.7), Inches(9), Inches(0.4),
         "注文完了画面が表示され、メールで注文確認が届きます", font_size=13, color=MUTED)

add_screenshot_placeholder(s, Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.3),
                           "注文確認画面")
add_screenshot_placeholder(s, Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.3),
                           "注文完了画面")
add_slide_num(s, 13, TOTAL)

# ================================================================
# SLIDE 14 — Payment Options
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_text(s, Inches(0), Inches(0.6), SLIDE_W, Inches(0.5),
         "選べる決済方法", font_size=30, bold=True, color=MAIN_DARK,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF)

pay_cards = [
    ("クレジットカード", "注文時にカード情報を入力。\nStripe社の安全な決済基盤を利用。",
     "個人・法人 どちらも利用可", SOFT, None),
    ("掛売（請求書払い）", "月末締め・翌月末払い\n飲食店の実務に合った決済です。",
     "法人のお客様向け", WHITE, MAIN),
    ("代金引換", "商品のお届け時に\n配送員へお支払い。",
     "個人・法人 どちらも利用可", SOFT, None),
]

pc_w = Inches(3.5)
pc_h = Inches(4.0)
pc_gap = Inches(0.35)
pc_total = pc_w * 3 + pc_gap * 2
pc_start = (SLIDE_W - pc_total) // 2
pc_y = Inches(1.6)

for i, (title, desc, badge, bg_col, border) in enumerate(pay_cards):
    x = Emu(int(pc_start) + int(Emu(int(pc_w) + int(pc_gap))) * i)
    if border:
        add_rounded_rect(s, x, pc_y, pc_w, pc_h, bg_col, border)
        star_y = pc_y - Inches(0.15)
        star_rect = add_rounded_rect(s, x + Inches(1.0), star_y, Inches(1.5), Inches(0.3), MAIN)
        add_text(s, x + Inches(1.0), star_y + Inches(0.02), Inches(1.5), Inches(0.28),
                 "おすすめ", font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    else:
        add_rounded_rect(s, x, pc_y, pc_w, pc_h, bg_col, RGBColor(0xE5, 0xE7, 0xEB))
    add_text(s, x + Inches(0.3), pc_y + Inches(0.5), pc_w - Inches(0.6), Inches(0.4),
             title, font_size=18, bold=True, color=MAIN_DARK, align=PP_ALIGN.CENTER)
    add_multiline(s, x + Inches(0.3), pc_y + Inches(1.2), pc_w - Inches(0.6), Inches(1.0),
                  desc.split('\n'), font_size=13, color=MUTED, align=PP_ALIGN.CENTER,
                  line_spacing=1.7)
    badge_rect = add_rounded_rect(s, x + Inches(0.5), pc_y + Inches(3.2),
                                   pc_w - Inches(1.0), Inches(0.35),
                                   HIGHLIGHT if not border else ACCENT)
    add_text(s, x + Inches(0.5), pc_y + Inches(3.22), pc_w - Inches(1.0), Inches(0.35),
             badge, font_size=11, bold=True, color=MAIN_DARK, align=PP_ALIGN.CENTER)

add_slide_num(s, 14, TOTAL)

# ================================================================
# SLIDE 15 — Invoice Payment Detail
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_rounded_rect(s, Inches(0.8), Inches(0.6), Inches(1.5), Inches(0.35), MAIN)
add_text(s, Inches(0.8), Inches(0.62), Inches(1.5), Inches(0.35),
         "法人向け", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
         "掛売（請求書払い）の流れ", font_size=28, bold=True, color=MAIN_DARK,
         font_name=FONT_SERIF)

invoice_steps = [
    "支払方法で「掛売」を選択して注文",
    "月中に複数回の注文もOK",
    "月末締めで請求書が発行される",
    "請求書一覧からPDF確認・ダウンロード",
    "翌月末までに銀行振込でお支払い",
]
for i, step in enumerate(invoice_steps):
    y = Inches(2.2) + Inches(i * 0.85)
    add_circle(s, Inches(1.2), y, Inches(0.45), MAIN)
    add_text(s, Inches(1.2), y + Inches(0.08), Inches(0.45), Inches(0.35),
             str(i + 1), font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.9), y + Inches(0.08), Inches(4.5), Inches(0.35),
             step, font_size=14, color=TEXT_COLOR)

add_screenshot_placeholder(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0),
                           "請求書一覧画面")
add_slide_num(s, 15, TOTAL)

# ================================================================
# SLIDE 16 — Delivery Status
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_rounded_rect(s, Inches(0.8), Inches(0.6), Inches(1.2), Inches(0.35), MAIN)
add_text(s, Inches(0.8), Inches(0.62), Inches(1.2), Inches(0.35),
         "注文後", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
         "配送状況の確認", font_size=28, bold=True, color=MAIN_DARK, font_name=FONT_SERIF)
add_text(s, Inches(0.8), Inches(1.7), Inches(9), Inches(0.4),
         "ダッシュボードから注文のステータスをリアルタイムで確認できます",
         font_size=13, color=MUTED)

add_screenshot_placeholder(s, Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.3),
                           "注文詳細 / 配送ステータス画面")

statuses = [
    ("出荷準備中", "商品を準備しています"),
    ("お届け中",   "配送業者に渡しました"),
    ("配達完了",   "お届けが完了しました"),
]

add_rounded_rect(s, Inches(7.0), Inches(2.4), Inches(5.5), Inches(2.5), SOFT,
                 RGBColor(0xE5, 0xE7, 0xEB))
add_text(s, Inches(7.2), Inches(2.6), Inches(5.0), Inches(0.35),
         "配送ステータスの見方", font_size=14, bold=True, color=MAIN_DARK)

for i, (status, desc) in enumerate(statuses):
    y = Inches(3.15) + Inches(i * 0.55)
    add_rounded_rect(s, Inches(7.4), y, Inches(1.6), Inches(0.35), HIGHLIGHT)
    add_text(s, Inches(7.4), y + Inches(0.03), Inches(1.6), Inches(0.3),
             status, font_size=12, bold=True, color=MAIN_DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.2), y + Inches(0.03), Inches(3.0), Inches(0.3),
             desc, font_size=12, color=MUTED)

add_rounded_rect(s, Inches(7.0), Inches(5.3), Inches(5.5), Inches(0.8), HIGHLIGHT)
add_text(s, Inches(7.2), Inches(5.4), Inches(5.1), Inches(0.6),
         "注文が入るとメール通知が届き、\nダッシュボードからもいつでも確認できます。",
         font_size=12, color=MAIN_DARK, align=PP_ALIGN.CENTER, bold=True)

add_slide_num(s, 16, TOTAL)

# ================================================================
# SLIDE 17 — Farm Pickup
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)

s.shapes.add_picture(img("farm_pickup.png"),
                     Inches(0), Inches(0), SLIDE_W, SLIDE_H)
add_bg_rect(s, WHITE, alpha=55)

add_rounded_rect(s, Inches(0.8), Inches(0.6), Inches(2.8), Inches(0.35), MAIN)
add_text(s, Inches(0.8), Inches(0.62), Inches(2.8), Inches(0.35),
         "セッツマルシェならでは", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
         "畑での受け取り", font_size=30, bold=True, color=MAIN_DARK, font_name=FONT_SERIF)

add_rounded_rect(s, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.5),
                 WHITE, RGBColor(0xE5, 0xE7, 0xEB))

pickup_steps = [
    ("チェックアウトで\n「畑で受け取り」を選択", "配送先の入力は不要"),
    ("出品者から\n準備完了の連絡", "受け取り日時を調整"),
    ("畑で受け取り", "生産者と直接会える体験"),
]

for i, (title, sub) in enumerate(pickup_steps):
    y = Inches(1.6) + Inches(i * 1.4)
    add_circle(s, Inches(7.4), y, Inches(0.5), MAIN)
    add_text(s, Inches(7.4), y + Inches(0.08), Inches(0.5), Inches(0.4),
             str(i + 1), font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.1), y, Inches(4.0), Inches(0.7),
             title, font_size=14, bold=True, color=MAIN_DARK, line_spacing=1.4)
    add_text(s, Inches(8.1), y + Inches(0.7), Inches(4.0), Inches(0.3),
             sub, font_size=11, color=MUTED)

add_rounded_rect(s, Inches(7.2), Inches(5.7), Inches(5.1), Inches(0.7), HIGHLIGHT)
add_text(s, Inches(7.4), Inches(5.75), Inches(4.7), Inches(0.55),
         "生産者の顔が見える安心感と、\n配送コストゼロのお得さ。",
         font_size=13, color=MAIN_DARK, align=PP_ALIGN.CENTER, bold=True)

add_bottom_stripe(s)
add_slide_num(s, 17, TOTAL)

# ================================================================
# SLIDE 18 — Convenient Features (Reorder + Templates)
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_text(s, Inches(0), Inches(0.6), SLIDE_W, Inches(0.5),
         "便利な機能", font_size=30, bold=True, color=MAIN_DARK,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF)
add_text(s, Inches(0), Inches(1.2), SLIDE_W, Inches(0.4),
         "毎日の仕入れをもっとラクに", font_size=15, color=MUTED,
         align=PP_ALIGN.CENTER)

feat_w = Inches(5.8)
feat_h = Inches(4.2)
feat_gap = Inches(0.4)
feat_total = feat_w * 2 + feat_gap
feat_start = (SLIDE_W - feat_total) // 2
feat_y = Inches(2.0)

features = [
    ("リピート注文", "過去の注文から「この注文をもう一度」\nボタンひとつで、同じ商品を\nカートに追加。\n毎回探す手間がありません。",
     "注文履歴から1クリック", "リピート注文ボタン"),
    ("注文テンプレート", "よく注文する商品の組み合わせを\nテンプレート保存。\n「いつもの」をワンクリックで\n注文できます。",
     "カートから保存も可能", "注文テンプレート一覧"),
]

for i, (title, desc, badge, ss_label) in enumerate(features):
    x = Emu(int(feat_start) + int(Emu(int(feat_w) + int(feat_gap))) * i)
    add_rounded_rect(s, x, feat_y, feat_w, feat_h, SOFT, RGBColor(0xE5, 0xE7, 0xEB))
    add_text(s, x + Inches(0.4), feat_y + Inches(0.3), feat_w - Inches(0.8), Inches(0.4),
             title, font_size=20, bold=True, color=MAIN_DARK)
    add_multiline(s, x + Inches(0.4), feat_y + Inches(0.9), feat_w - Inches(0.8), Inches(1.8),
                  desc.split('\n'), font_size=12, color=MUTED, line_spacing=1.7)
    add_rounded_rect(s, x + Inches(0.4), feat_y + Inches(2.8), Inches(2.5), Inches(0.35), HIGHLIGHT)
    add_text(s, x + Inches(0.4), feat_y + Inches(2.82), Inches(2.5), Inches(0.35),
             badge, font_size=11, bold=True, color=MAIN_DARK, align=PP_ALIGN.CENTER)
    add_screenshot_placeholder(s, x + Inches(3.2), feat_y + Inches(2.5), Inches(2.2), Inches(1.4),
                               ss_label)

add_slide_num(s, 18, TOTAL)

# ================================================================
# SLIDE 19 — B2B Features (Pricing, Approval, Org)
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg_rect(s, WHITE)
add_bottom_stripe(s)

add_text(s, Inches(0), Inches(0.6), SLIDE_W, Inches(0.5),
         "法人向けの安心機能", font_size=30, bold=True, color=MAIN_DARK,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF)
add_text(s, Inches(0), Inches(1.2), SLIDE_W, Inches(0.4),
         "組織での利用に対応した仕組みをご用意しています", font_size=15, color=MUTED,
         align=PP_ALIGN.CENTER)

b2b_cards = [
    ("顧客別価格", "お取引状況に応じた\n特別価格を設定。\nログインすると自動で適用。", MAIN),
    ("購買承認フロー", "担当者が注文 → 店長が承認。\n組織に合わせた承認\nステップを設定可能。", ACCENT),
    ("組織メンバー管理", "注文担当・承認者・経理など\nメンバーごとに\n権限を設定できます。", MAIN_DARK),
]

bc_w = Inches(3.5)
bc_h = Inches(4.0)
bc_gap = Inches(0.35)
bc_total = bc_w * 3 + bc_gap * 2
bc_start = (SLIDE_W - bc_total) // 2
bc_y = Inches(2.0)

for i, (title, desc, bar_color) in enumerate(b2b_cards):
    x = Emu(int(bc_start) + int(Emu(int(bc_w) + int(bc_gap))) * i)
    add_rounded_rect(s, x, bc_y, bc_w, bc_h, SOFT, RGBColor(0xE5, 0xE7, 0xEB))
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, bc_y, bc_w, Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = bar_color
    top_bar.line.fill.background()
    add_circle(s, x + (bc_w - Inches(0.7)) // 2, bc_y + Inches(0.4), Inches(0.7), bar_color)
    icons = ["¥", "✓", "⚙"]
    add_text(s, x + (bc_w - Inches(0.7)) // 2, bc_y + Inches(0.48), Inches(0.7), Inches(0.55),
             icons[i], font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), bc_y + Inches(1.4), bc_w - Inches(0.6), Inches(0.4),
             title, font_size=18, bold=True, color=MAIN_DARK, align=PP_ALIGN.CENTER)
    add_multiline(s, x + Inches(0.3), bc_y + Inches(2.0), bc_w - Inches(0.6), Inches(1.5),
                  desc.split('\n'), font_size=12, color=MUTED, align=PP_ALIGN.CENTER,
                  line_spacing=1.7)

add_slide_num(s, 19, TOTAL)

# ================================================================
# SLIDE 20 — Closing
# ================================================================
s = prs.slides.add_slide(BLANK)
add_gradient_bg(s, MAIN, MAIN_DARK)

s.shapes.add_picture(img("circulation_vision.png"),
                     Inches(0), Inches(0), SLIDE_W, SLIDE_H)
add_bg_rect(s, MAIN_DARK, alpha=75)

add_text(s, Inches(0), Inches(1.8), SLIDE_W, Inches(1.8),
         "地域でつくり、地域で食べる。\nその循環を、ごいっしょに。",
         font_size=34, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name=FONT_SERIF, line_spacing=1.6)

add_text(s, Inches(2.5), Inches(3.8), Inches(8.3), Inches(0.7),
         "セッツマルシェは、飲食店のみなさまと一緒に育てていくサービスです。\nご意見・ご要望をぜひお聞かせください。",
         font_size=15, color=RGBColor(0xCC, 0xCC, 0xCC),
         align=PP_ALIGN.CENTER, line_spacing=1.7)

add_rounded_rect(s, Inches(3.8), Inches(5.0), Inches(5.7), Inches(1.2), WHITE)

add_text(s, Inches(3.8), Inches(5.15), Inches(5.7), Inches(0.9),
         "セッツマルシェ（渡辺ファーム）\nお問い合わせ：サイト内「お問い合わせフォーム」より",
         font_size=13, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.7)

# ── save ───────────────────────────────────────────
prs.save(OUT)
print(f"✅ Saved: {OUT}")
print(f"   Slides: {len(prs.slides)}")
